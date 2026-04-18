import os
import shutil
import ast
import json
from datetime import datetime
from typing import List, Optional
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from dotenv import load_dotenv
from contextlib import asynccontextmanager

# Load environment variables
load_dotenv()

# Firebase + repo (imports early so Admin SDK picks up env vars)
from firebase_admin_init import db as _fb_db  # noqa: F401 — init side effect
import firestore_repo as repo
from auth import get_current_user, optional_current_user, require_teacher
from analytics import router as analytics_router

# Resolve API key once so all clients use the same credential source.
RESOLVED_GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
if not RESOLVED_GEMINI_API_KEY:
    print("WARNING: GEMINI_API_KEY/GOOGLE_API_KEY not found in environment variables.")

# LlamaIndex Imports
from llama_index.core import (
    VectorStoreIndex,
    SimpleDirectoryReader,
    StorageContext,
    Settings,
    Document
)
from llama_index.core.node_parser import SemanticSplitterNodeParser, CodeSplitter, TokenTextSplitter
from llama_index.core.tools import QueryEngineTool
from llama_index.core.query_engine import RouterQueryEngine
from llama_index.core.selectors import LLMSingleSelector
from llama_index.llms.google_genai import GoogleGenAI
from llama_index.embeddings.google_genai import GoogleGenAIEmbedding
from llama_index.vector_stores.chroma import ChromaVectorStore
import chromadb

# Google GenAI SDK (for Vision/Digitize)
from google import genai
from google.genai import types
import logging

# File Generation Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from xhtml2pdf import pisa
import markdown
import re
import random
import uuid
import requests
import base64
import io

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# --- Configuration ---
PERSIST_DIR = "./chroma_db"
UPLOAD_DIR = "./uploads"
METADATA_FILE = os.path.join(PERSIST_DIR, "materials_metadata.json")

# --- Setup Models ---
# Use a single resolved key and provide clearer startup failures.
if RESOLVED_GEMINI_API_KEY:
    os.environ["GOOGLE_API_KEY"] = RESOLVED_GEMINI_API_KEY

try:
    Settings.llm = GoogleGenAI(
        model="models/gemini-2.0-flash",
        api_key=RESOLVED_GEMINI_API_KEY
    )
    # NOTE: the llama_index google_genai package uses the kwarg `model_name`,
    # not `model`. Passing `model=` silently falls through to the default
    # `gemini-embedding-2-preview`, which is rate-limited and returns empty
    # vectors on short inputs. Pin to the stable GA `gemini-embedding-001`.
    Settings.embed_model = GoogleGenAIEmbedding(
        model_name="gemini-embedding-001",
        api_key=RESOLVED_GEMINI_API_KEY,
    )
except Exception as e:
    raise RuntimeError(
        "Failed to initialize Gemini models. Check network/DNS/proxy settings and ensure "
        "you can resolve and reach Google's API endpoints (for example, generativelanguage.googleapis.com)."
    ) from e

# --- Helper Section: Code Structure Analysis ---
def extract_code_structure(code: str, language: str) -> dict:
    """
    Extracts structural information from code (functions, classes, imports, etc.)
    for syntax-aware search.
    """
    structure = {
        "functions": [],
        "classes": [],
        "imports": [],
        "variables": [],
        "patterns": []
    }
    
    if language in ['python', 'py']:
        try:
            tree = ast.parse(code)
            
            # Extract functions
            for node in ast.walk(tree):
                if isinstance(node, ast.FunctionDef):
                    func_info = {
                        "name": node.name,
                        "line": node.lineno,
                        "args": [arg.arg for arg in node.args.args],
                        "docstring": ast.get_docstring(node) or ""
                    }
                    structure["functions"].append(func_info)
                    
                    # Detect patterns (recursion, loops, etc.)
                    # Check for recursive calls (function calling itself)
                    for child in ast.walk(node):
                        if isinstance(child, ast.Call) and isinstance(child.func, ast.Name):
                            if child.func.id == node.name:
                                structure["patterns"].append(f"recursive_function:{node.name}")
                                break
                    
                elif isinstance(node, ast.ClassDef):
                    class_info = {
                        "name": node.name,
                        "line": node.lineno,
                        "methods": [n.name for n in node.body if isinstance(n, ast.FunctionDef)],
                        "docstring": ast.get_docstring(node) or ""
                    }
                    structure["classes"].append(class_info)
                    
                elif isinstance(node, (ast.Import, ast.ImportFrom)):
                    if isinstance(node, ast.Import):
                        imports = [alias.name for alias in node.names]
                    else:
                        imports = [f"{node.module}.{alias.name}" if node.module else alias.name 
                                 for alias in node.names]
                    structure["imports"].extend(imports)
                    
        except SyntaxError:
            # If code has syntax errors, try regex fallback
            pass
    
    # Regex-based extraction for other languages
    if language in ['c', 'cpp', 'c++']:
        # Extract function definitions
        func_pattern = r'(?:^|\n)\s*(?:\w+\s+)*(\w+)\s*\([^)]*\)\s*\{'
        functions = re.findall(func_pattern, code, re.MULTILINE)
        structure["functions"] = [{"name": f, "line": 0} for f in functions]
        
        # Extract includes
        includes = re.findall(r'#include\s+[<"]([^>"]+)[>"]', code)
        structure["imports"] = includes
        
    elif language in ['java']:
        # Extract class definitions
        class_pattern = r'(?:public\s+)?(?:class|interface)\s+(\w+)'
        classes = re.findall(class_pattern, code)
        structure["classes"] = [{"name": c, "line": 0} for c in classes]
        
        # Extract method definitions
        method_pattern = r'(?:public|private|protected)?\s*(?:\w+\s+)*(\w+)\s*\([^)]*\)'
        methods = re.findall(method_pattern, code)
        structure["functions"] = [{"name": m, "line": 0} for m in methods]
        
        # Extract imports
        imports = re.findall(r'import\s+([\w.]+);', code)
        structure["imports"] = imports
        
    elif language in ['javascript', 'js', 'typescript', 'ts']:
        # Extract function/arrow function definitions
        func_pattern = r'(?:function\s+(\w+)|const\s+(\w+)\s*=\s*(?:\([^)]*\)\s*=>|function))'
        matches = re.findall(func_pattern, code)
        functions = [m[0] or m[1] for m in matches if m[0] or m[1]]
        structure["functions"] = [{"name": f, "line": 0} for f in functions]
        
        # Extract class definitions
        class_pattern = r'class\s+(\w+)'
        classes = re.findall(class_pattern, code)
        structure["classes"] = [{"name": c, "line": 0} for c in classes]
        
        # Extract imports
        imports = re.findall(r'import\s+(?:.*\s+from\s+)?[\'"]([^\'"]+)[\'"]', code)
        structure["imports"] = imports
    
    return structure


def split_code_by_functions(code: str, language: str, metadata: dict) -> List[Document]:
    """
    Splits code into function-level chunks for better semantic search.
    Each function/class becomes its own searchable chunk with proper metadata.
    Also creates a header chunk for imports and module-level code.
    """
    chunks = []
    lines = code.split('\n')
    
    if language in ['python', 'py']:
        try:
            tree = ast.parse(code)
            
            # Find all top-level functions and classes with their line ranges
            items = []
            for node in ast.iter_child_nodes(tree):
                if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
                    # Find the end line (last line of the function)
                    end_line = max(
                        getattr(child, 'end_lineno', node.lineno)
                        for child in ast.walk(node)
                        if hasattr(child, 'end_lineno')
                    )
                    items.append({
                        'type': 'function',
                        'name': node.name,
                        'start': node.lineno - 1,  # Convert to 0-indexed
                        'end': end_line,
                        'docstring': ast.get_docstring(node) or "",
                        'args': [arg.arg for arg in node.args.args]
                    })
                elif isinstance(node, ast.ClassDef):
                    end_line = max(
                        getattr(child, 'end_lineno', node.lineno)
                        for child in ast.walk(node)
                        if hasattr(child, 'end_lineno')
                    )
                    methods = [n.name for n in node.body if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))]
                    items.append({
                        'type': 'class',
                        'name': node.name,
                        'start': node.lineno - 1,
                        'end': end_line,
                        'docstring': ast.get_docstring(node) or "",
                        'methods': methods
                    })
            
            # Sort by start line
            items.sort(key=lambda x: x['start'])
            
            # Create header chunk (imports, module docstring, etc.)
            if items:
                header_end = items[0]['start']
                if header_end > 0:
                    header_content = '\n'.join(lines[:header_end])
                    if header_content.strip():
                        header_doc = Document(
                            text=header_content,
                            metadata={
                                **metadata,
                                'chunk_type': 'header',
                                'chunk_name': 'module_header',
                                'contains_imports': 'True'
                            }
                        )
                        chunks.append(header_doc)
            
            # Create a chunk for each function/class
            for item in items:
                func_content = '\n'.join(lines[item['start']:item['end']])
                
                # Create rich metadata for this function/class
                chunk_metadata = {
                    **metadata,
                    'chunk_type': item['type'],
                    'chunk_name': item['name'],
                    'function_names': item['name'] if item['type'] == 'function' else '',
                    'class_names': item['name'] if item['type'] == 'class' else '',
                    'docstring_preview': item['docstring'][:200] if item['docstring'] else '',
                }
                
                if item['type'] == 'function':
                    chunk_metadata['function_args'] = ', '.join(item.get('args', []))
                elif item['type'] == 'class':
                    chunk_metadata['class_methods'] = ', '.join(item.get('methods', []))
                
                func_doc = Document(
                    text=func_content,
                    metadata=chunk_metadata
                )
                chunks.append(func_doc)
            
            # If there's code after the last function (like if __name__ == "__main__")
            if items:
                last_end = items[-1]['end']
                if last_end < len(lines):
                    footer_content = '\n'.join(lines[last_end:])
                    if footer_content.strip():
                        footer_doc = Document(
                            text=footer_content,
                            metadata={
                                **metadata,
                                'chunk_type': 'footer',
                                'chunk_name': 'main_execution'
                            }
                        )
                        chunks.append(footer_doc)
            
            if chunks:
                print(f"[DEBUG] Function-aware splitter created {len(chunks)} chunks")
                return chunks
                
        except SyntaxError as e:
            print(f"[DEBUG] AST parsing failed: {e}, falling back to regex")
    
    # Fallback: Use regex for other languages or if AST fails
    if language in ['python', 'py']:
        # Regex fallback for Python
        func_pattern = r'^(def\s+\w+|class\s+\w+)'
    elif language in ['javascript', 'js', 'typescript', 'ts']:
        func_pattern = r'^(function\s+\w+|class\s+\w+|const\s+\w+\s*=\s*(?:\([^)]*\)\s*=>|function))'
    elif language in ['c', 'cpp', 'c++']:
        func_pattern = r'^(?:\w+\s+)+(\w+)\s*\([^)]*\)\s*\{'
    elif language in ['java']:
        func_pattern = r'^(?:public|private|protected)?\s*(?:static\s+)?(?:\w+\s+)+(\w+)\s*\([^)]*\)'
    else:
        # Generic fallback - just use the whole document
        return []
    
    # Simple line-based splitting if no function boundaries found
    return []


def is_code_query(query: str) -> bool:
    """
    Detects if a query is code-specific (looking for code patterns, functions, etc.)
    """
    query_lower = query.lower()
    
    # Code-specific indicators
    code_indicators = [
        "function", "class", "method", "import", "define", "implement",
        "code for", "snippet", "example", "how to write", "how to create",
        "recursive", "algorithm", "data structure", "pattern", "syntax",
        "api", "library", "module", "package"
    ]
    
    # Function/class name patterns
    if re.search(r'\b(function|def|class|interface|method)\s+\w+', query_lower):
        return True
    
    # Code pattern queries
    if any(indicator in query_lower for indicator in code_indicators):
        return True
    
    return False

def normalize_week(week_str: str) -> str:
    """
    Normalizes week string to a consistent format for comparison.
    Handles: "7", "Week 7", "week 7", "WEEK 7", etc.
    Returns normalized format: "Week 7" or just the number if no "week" prefix.
    """
    if not week_str:
        return ""
    week_str = str(week_str).strip()
    # Extract number
    match = re.search(r'(\d+)', week_str)
    if match:
        num = match.group(1)
        # If it already has "week" prefix, keep it; otherwise add it
        if 'week' in week_str.lower():
            return f"Week {num}"
        return num
    return week_str

def extract_metadata_filters(query: str) -> dict:
    """
    Extracts metadata filters from natural language queries.
    Returns dict with week, topic, course, tags, category filters.
    """
    query_lower = query.lower()
    filters = {}
    
    # Extract week number (normalized)
    week_patterns = [
        r'week\s*(\d+)',
        r'week\s*(\d+)\s*content',
        r'week\s*(\d+)\s*materials',
        r'lecture\s*(\d+)',
        r'lecture\s*(\d+)\s*content',
    ]
    for pattern in week_patterns:
        match = re.search(pattern, query_lower)
        if match:
            week_num = match.group(1)
            # Store both formats for flexible matching
            filters['week'] = week_num  # Store as number for flexible matching
            filters['week_normalized'] = f"Week {week_num}"  # Also store normalized
            break
    
    # Extract course name (common patterns)
    course_patterns = [
        r'(?:course|subject|class)\s+(?:called|named|is|about)\s+([a-z\s]+?)(?:\s|$|,|\.)',
        r'for\s+([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)\s+course',
    ]
    for pattern in course_patterns:
        match = re.search(pattern, query_lower)
        if match:
            filters['course'] = match.group(1).strip()
            break
    
    # Extract topic (look for "about X", "on X", "regarding X")
    topic_patterns = [
        r'(?:about|on|regarding|concerning|related to)\s+([a-z\s]+?)(?:\s|$|,|\.|\?)',
        r'topic\s+(?:is|of|about)\s+([a-z\s]+?)(?:\s|$|,|\.)',
    ]
    for pattern in topic_patterns:
        match = re.search(pattern, query_lower)
        if match:
            topic = match.group(1).strip()
            # Filter out common stop words
            if topic and len(topic) > 2 and topic not in ['the', 'a', 'an', 'is', 'are']:
                filters['topic'] = topic
                break
    
    # Extract category (theory/lab)
    if any(word in query_lower for word in ['theory', 'theoretical', 'lecture', 'reading', 'notes', 'slides']):
        filters['category'] = 'theory'
    elif any(word in query_lower for word in ['lab', 'laboratory', 'code', 'programming', 'practical', 'exercise']):
        filters['category'] = 'lab'
    
    # Extract tags (words after # or common tag indicators)
    tag_patterns = [
        r'#(\w+)',
        r'tag(?:ged|s)?\s+(?:as|with|:)?\s*([a-z\s,]+)',
    ]
    for pattern in tag_patterns:
        matches = re.findall(pattern, query_lower)
        if matches:
            tags = []
            for match in matches:
                if isinstance(match, tuple):
                    tags.extend([t.strip() for t in match if t.strip()])
                else:
                    tags.extend([t.strip() for t in match.split(',') if t.strip()])
            if tags:
                filters['tags'] = tags
                break
    
    return filters

def extract_code_entities(query: str) -> dict:
    """
    Extracts code-related entities from a query (function names, class names, etc.)
    """
    entities = {
        "function_names": [],
        "class_names": [],
        "imports": [],
        "patterns": []
    }
    
    # Extract potential function/class names (capitalized or common patterns)
    words = re.findall(r'\b[A-Z][a-zA-Z0-9_]*\b', query)
    entities["class_names"] = words
    
    # Extract lowercase identifiers that might be functions
    func_candidates = re.findall(r'\b[a-z_][a-zA-Z0-9_]*\b', query)
    # Filter out common words
    common_words = {'the', 'a', 'an', 'is', 'are', 'was', 'were', 'how', 'what', 'where', 'when', 'why', 'does', 'do', 'can', 'should', 'will', 'find', 'search', 'show', 'get'}
    entities["function_names"] = [w for w in func_candidates if w not in common_words and len(w) > 2]
    
    # Detect pattern queries
    if 'recursive' in query.lower() or 'recursion' in query.lower():
        entities["patterns"].append("recursive")
    if 'loop' in query.lower() or 'iteration' in query.lower():
        entities["patterns"].append("loop")
    if 'sort' in query.lower():
        entities["patterns"].append("sort")
    if 'search' in query.lower() and 'algorithm' in query.lower():
        entities["patterns"].append("search_algorithm")
    
    return entities

# --- Helper Section ---
def get_language_from_ext(ext: str) -> str:
    """Map file extension to tree-sitter language."""
    ext = ext.lower().strip()
    mapping = {
        ".js": "javascript", ".jsx": "javascript", ".mjs": "javascript",
        ".ts": "typescript", ".tsx": "typescript",
        ".py": "python",
        ".java": "java",
        ".c": "c", ".h": "c",
        ".cpp": "cpp", ".cc": "cpp", ".cxx": "cpp", ".hpp": "cpp",
        ".cs": "c_sharp",
        ".go": "go",
        ".rs": "rust",
        ".html": "html", ".htm": "html",
        ".css": "css",
        ".json": "json"
    }
    return mapping.get(ext, "python")

def validate_code_syntax(code_str: str) -> bool:
    try:
        ast.parse(code_str)
        return True
    except SyntaxError:
        return False

def load_metadata() -> dict:
    """Return all materials keyed by file_id. Backed by Firestore."""
    try:
        return {m["id"]: m for m in repo.list_materials()}
    except Exception as e:
        logger.error(f"Failed to load materials from Firestore: {e}")
        return {}


def save_metadata(data: dict):
    """Deprecated: retained as a no-op shim. Use repo.upsert_material / repo.delete_material directly."""
    logger.warning(
        "save_metadata() is a no-op shim; update callers to use firestore_repo directly."
    )

# --- Helper: Code Detection & Cleaning ---
def is_code_request(prompt: str) -> bool:
    """
    Heuristic to detect if a prompt is asking for code generation vs a conceptual explanation.
    """
    prompt_lower = prompt.lower().strip()
    
    # Code request indicators (action verbs)
    code_indicators = [
        "write", "implement", "create", "generate", "code for",
        "build", "make a", "develop", "program", "script",
        "function for", "algorithm for", "solve", "calculate"
    ]
    
    # Conceptual question indicators
    concept_indicators = [
        "what is", "explain", "how does", "why is", "define",
        "describe", "what are", "difference between", "compare",
        "introduction to", "overview of", "tell me about"
    ]
    
    # Check for code indicators first
    for indicator in code_indicators:
        if indicator in prompt_lower:
            return True
    
    # Check for concept indicators
    for indicator in concept_indicators:
        if prompt_lower.startswith(indicator) or f" {indicator}" in prompt_lower:
            return False
    
    # Default: If it ends with a question mark, likely conceptual   
    if "?" in prompt_lower or "how to" in prompt_lower:
        return False
    
    # Additional specific check: "how to" often means explanation + code
    if prompt_lower.startswith("how to"):
         return False

    # Default to code if unsure
    return True

def detect_language_from_prompt(prompt: str) -> tuple[str, str]:
    """
    Detect programming language from the prompt using regex for accuracy.
    Returns (language_name, file_extension).
    """
    prompt_lower = prompt.lower()
    
    # Order matters: check C++ before C
    language_patterns = [
        (r'\b(python|py)\b', "python", ".py"),
        (r'\b(c\+\+|cpp)\b', "cpp", ".cpp"),
        (r'\b(java)\b', "java", ".java"),
        (r'\b(javascript|js)\b', "javascript", ".js"),
        (r'\b(typescript|ts)\b', "typescript", ".ts"),
        (r'\b(c)\b', "c", ".c"), # Matches isolated 'c'
        (r'\b(c#|csharp)\b', "csharp", ".cs"),
        (r'\b(go|golang)\b', "go", ".go"),
        (r'\b(rust)\b', "rust", ".rs"),
        (r'\b(ruby|rb)\b', "ruby", ".rb"),
        (r'\b(php)\b', "php", ".php"),
        (r'\b(html)\b', "html", ".html"),
        (r'\b(css)\b', "css", ".css"),
        (r'\b(sql)\b', "sql", ".sql"),
        (r'\b(bash|shell)\b', "bash", ".sh"),
    ]
    
    for pattern, lang, ext in language_patterns:
        if re.search(pattern, prompt_lower):
            return (lang, ext)
    
    # Default to Python if nothing found
    return ("python", ".py")

def clean_code_content(content: str, lang_name: str) -> str:
    """
    Cleans common LLM artifacts from code generation.
    """
    lines = content.strip().split('\n')
    if not lines: return content
    
    first_line = lines[0].strip().lower()
    
    # Check if first line is just the language name
    possible_names = [lang_name, "c", "cpp", "c++", "python", "java", "code", "block"]
    
    if first_line in possible_names or first_line.replace("#", "").strip() in possible_names:
        # Check if it's NOT a valid code line (e.g. #include)
        if not first_line.startswith("#include") and not first_line.startswith("import"):
             return "\n".join(lines[1:]).strip()
             
    return content

# --- Helper: Code Validation ---
def validate_code_syntax(code_str: str) -> bool:
    """Basic Python syntax validation using AST."""
    try:
        ast.parse(code_str)
        return True
    except SyntaxError:
        return False

def validate_locally(code: str, language: str) -> dict:
    """
    Tiered Local Validation (Fallback)
    Tier 2: Static Analysis (Python AST)
    Tier 3: Heuristics (Brackets)
    Tier 4: Security Keywords
    """
    # Tier 4: Security
    dangerous_keywords = ['system', 'fork', 'exec', 'popen', 'subprocess', 'os.system']
    for kw in dangerous_keywords:
        if re.search(f'\\b{kw}\\b', code):
             return {"valid": False, "output": f"Security Check Failed: Forbidden keyword '{kw}' detected.", "error": "Security Risk"}

    # Tier 2: Static Analysis (Python)
    if language in ['python', 'py']:
        try:
            ast.parse(code)
            return {"valid": True, "output": "Local Validation: Python Syntax Valid (AST Checked).", "error": None}
        except SyntaxError as e:
            return {"valid": False, "output": f"Local Validation Error: {e}", "error": str(e)}

    # Tier 3: Heuristics (C/C++/Java/JS) - Brackets & structure
    # 1. Bracket Balance
    brackets = {'(': ')', '{': '}', '[': ']'}
    stack = []
    for char in code:
        if char in brackets.keys():
            stack.append(char)
        elif char in brackets.values():
            if not stack:
                return {"valid": False, "output": "Local Validation Error: Unbalanced brackets (extra closing).", "error": "Unbalanced Brackets"}
            last = stack.pop()
            if brackets[last] != char:
                return {"valid": False, "output": "Local Validation Error: Mismatched brackets.", "error": "Mismatched Brackets"}
    if stack:
        return {"valid": False, "output": "Local Validation Error: Unbalanced brackets (missing closing).", "error": "Unbalanced Brackets"}

    return {"valid": True, "output": "Local Validation: Structure looks correct (Heuristic Check).", "error": None}

def validate_code_with_fallback(code: str, language: str) -> dict:
    """
    Validates code using Wandbox -> Fallback to Local.
    """
    # Try Wandbox First
    url = "https://wandbox.org/api/compile.json"
    
    compiler_map = {
        "python": "cpython-3.14.0", "py": "cpython-3.14.0",
        "c": "gcc-13.2.0-c",
        "cpp": "gcc-13.2.0", "c++": "gcc-13.2.0",
        "java": "openjdk-jdk-22+36",
        "javascript": "nodejs-20.17.0", "js": "nodejs-20.17.0",
        "typescript": "typescript-5.6.2", "ts": "typescript-5.6.2",
        "go": "go-1.23.2",
        "rust": "rust-1.82.0",
        "ruby": "ruby-3.4.1",
        "swift": "swift-6.0.1",
        "php": "php-8.3.12",
        "sql": "sqlite-3.46.1"
    }
    
    compiler = compiler_map.get(language.lower())
    if compiler:
        payload = {
            "code": code,
            "compiler": compiler,
            "save": False
        }
        try:
            print(f"DEBUG: Validating {language} code with Wandbox...")
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.post(url, json=payload, headers=headers, timeout=10)
            
            # If API error (500 etc), raise exception to trigger fallback
            if response.status_code != 200:
                raise Exception(f"Wandbox API Error {response.status_code}")
                
            data = response.json()
            if data.get("status") == "0":
                output = data.get("program_output", "") + data.get("compiler_output", "")
                return {"valid": True, "output": output, "error": None}
            else:
                 # Real compilation error - Return strictly
                error_log = data.get("compiler_error") or data.get("program_error") or "Unknown Error"
                return {"valid": False, "output": error_log, "error": error_log}
                
        except Exception as e:
            print(f"DEBUG: Wandbox failed ({e}). Falling back to Local Validator.")
            # Fallback proceeds below...
    else:
        print(f"DEBUG: Language {language} not in Wandbox map. Using Local Validator.")

    # FALLBACK: Local Validation
    return validate_locally(code, language)

# --- Helpers: Image Generation ---
def generate_image(prompt: str, topic_context: str = "") -> Optional[str]:
    """
    Generates an image using Gemini image generation models.
    Returns the filepath to the saved image, or None if generation fails.
    """
    try:
        # Use Gemini image generation model
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
        
        # Enhanced prompt with context
        full_prompt = f"{prompt}"
        if topic_context:
            full_prompt = f"Create an educational diagram or visual aid for: {topic_context}. {prompt}. Make it clear, professional, and suitable for academic materials."
        
        print(f"DEBUG: Generating image with prompt: {full_prompt[:100]}...")
        
        # Try image generation models in order of preference
        image_models = [
            'models/gemini-2.5-flash-image',
            'models/gemini-2.0-flash-exp-image-generation',
            'models/gemini-3-pro-image-preview'
        ]
        
        image_data = None
        for model_name in image_models:
            try:
                # Do NOT use response_mime_type - Gemini image models return images in parts
                response = client.models.generate_content(
                    model=model_name,
                    contents=[full_prompt]
                )
                
                # Extract image from response parts
                # Response can have multiple parts - some text, some images
                # Try direct parts access first (newer API format)
                if hasattr(response, 'parts'):
                    for part in response.parts:
                        if hasattr(part, 'inline_data') and part.inline_data is not None:
                            # Try using as_image() method (recommended approach)
                            try:
                                if hasattr(part, 'as_image'):
                                    img = part.as_image()
                                    # Convert PIL Image to bytes
                                    img_bytes = io.BytesIO()
                                    # Try with format parameter, fallback without it
                                    try:
                                        img.save(img_bytes, format='PNG')
                                    except TypeError:
                                        # Some PIL versions or image types don't accept format parameter
                                        img.save(img_bytes, 'PNG')
                                    image_data = img_bytes.getvalue()
                                    print(f"DEBUG: Successfully extracted image using as_image()")
                                    break
                            except Exception as img_err:
                                print(f"DEBUG: as_image() failed ({img_err}), trying direct data access...")
                                # Fallback to direct data access
                                try:
                                    if hasattr(part.inline_data, 'data'):
                                        raw_data = part.inline_data.data
                                        # Handle both bytes and base64 string
                                        if isinstance(raw_data, bytes):
                                            image_data = raw_data
                                        elif isinstance(raw_data, str):
                                            image_data = base64.b64decode(raw_data)
                                        else:
                                            image_data = raw_data
                                        print(f"DEBUG: Successfully extracted image using inline_data.data")
                                        break
                                except Exception as data_err:
                                    print(f"DEBUG: Direct data access failed: {data_err}")
                                    continue
                
                # Also check candidates structure (alternative API format)
                if not image_data and hasattr(response, 'candidates') and response.candidates:
                    candidate = response.candidates[0]
                    if hasattr(candidate, 'content') and candidate.content:
                        if hasattr(candidate.content, 'parts'):
                            for part in candidate.content.parts:
                                if hasattr(part, 'inline_data') and part.inline_data is not None:
                                    try:
                                        if hasattr(part, 'as_image'):
                                            img = part.as_image()
                                            img_bytes = io.BytesIO()
                                            # Try with format parameter, fallback without it
                                            try:
                                                img.save(img_bytes, format='PNG')
                                            except TypeError:
                                                img.save(img_bytes, 'PNG')
                                            image_data = img_bytes.getvalue()
                                            print(f"DEBUG: Successfully extracted image from candidate using as_image()")
                                            break
                                        elif hasattr(part.inline_data, 'data'):
                                            raw_data = part.inline_data.data
                                            # Handle both bytes and base64 string
                                            if isinstance(raw_data, bytes):
                                                image_data = raw_data
                                            elif isinstance(raw_data, str):
                                                image_data = base64.b64decode(raw_data)
                                            else:
                                                image_data = raw_data
                                            print(f"DEBUG: Successfully extracted image from candidate using data")
                                            break
                                    except Exception as img_err:
                                        print(f"DEBUG: Error extracting from candidate: {img_err}")
                                        continue
                
                if image_data:
                    break
            except Exception as e:
                print(f"DEBUG: Model {model_name} failed: {e}")
                continue
        
        if image_data:
            # Decode base64 if needed (inline_data.data might be base64 string)
            if isinstance(image_data, str):
                try:
                    image_data = base64.b64decode(image_data)
                except Exception:
                    # If it's not base64, try encoding it
                    image_data = image_data.encode('latin-1') if isinstance(image_data, str) else image_data
            
            # Save image
            img_filename = f"img_{uuid.uuid4().hex[:8]}.png"
            img_filepath = os.path.join(UPLOAD_DIR, img_filename)
            
            # Ensure directory exists
            os.makedirs(UPLOAD_DIR, exist_ok=True)
            
            # Save the image
            with open(img_filepath, "wb") as img_file:
                if isinstance(image_data, bytes):
                    img_file.write(image_data)
                else:
                    # Convert to bytes if needed
                    img_file.write(bytes(image_data))
            
            print(f"DEBUG: Image saved to {img_filepath}")
            return img_filepath
        else:
            print("DEBUG: No image data received from generation")
            return None
            
    except Exception as e:
        print(f"DEBUG: Image generation failed: {e}")
        import traceback
        traceback.print_exc()
        return None

def extract_image_placeholders(content: str) -> List[dict]:
    """
    Extracts image placeholders from markdown content.
    Looks for patterns like [IMAGE: description] or ![alt](IMAGE:description)
    Returns list of dicts with 'placeholder', 'description', and 'position'
    """
    image_placeholders = []
    
    # Pattern 1: [IMAGE: description]
    pattern1 = r'\[IMAGE:\s*([^\]]+)\]'
    for match in re.finditer(pattern1, content):
        image_placeholders.append({
            'placeholder': match.group(0),
            'description': match.group(1).strip(),
            'position': match.start()
        })
    
    # Pattern 2: ![alt](IMAGE:description)
    pattern2 = r'!\[([^\]]*)\]\(IMAGE:\s*([^\)]+)\)'
    for match in re.finditer(pattern2, content):
        image_placeholders.append({
            'placeholder': match.group(0),
            'description': match.group(2).strip(),
            'position': match.start()
        })
    
    # Pattern 3: <!--IMAGE: description-->
    pattern3 = r'<!--IMAGE:\s*([^-]+)-->'
    for match in re.finditer(pattern3, content):
        image_placeholders.append({
            'placeholder': match.group(0),
            'description': match.group(1).strip(),
            'position': match.start()
        })
    
    # Sort by position
    image_placeholders.sort(key=lambda x: x['position'])
    return image_placeholders

# --- Helpers: File Generation ---
def create_pptx_from_text(content: str, filename: str, topic_context: str = "") -> str:
    """Creates a PPTX with proper image support using dedicated visual aid slides."""
    print(f"DEBUG: Creating PPTX for {filename}...")
    prs = Presentation()
    
    # CLEANUP: Remove conversational preambles from AI responses
    conversational_patterns = [
        r'^(?:Okay|Sure|Alright|Here are|I\'ll create|I will create|Here is|Let me)[^.!?]*(?:slides?|presentation|PPTX)[^.!?]*[.!?\n]+',
        r'^(?:Based on|Using|Adhering to|Following)[^.!?]*[.!?\n]+',
        r'^(?:Here\'s|This is)[^.!?]*(?:slides?|presentation)[^.!?]*[.!?\n]+',
    ]
    for pattern in conversational_patterns:
        content = re.sub(pattern, '', content, flags=re.IGNORECASE | re.MULTILINE)
    
    # Also strip any leading text before "Slide 1" or "## "
    if "Slide 1" in content:
        idx = content.find("Slide 1")
        if idx > 0 and idx < 500:  # Only strip if within first 500 chars
            content = content[idx:]
    elif "## " in content:
        idx = content.find("## ")
        if idx > 0 and idx < 500:
            content = content[idx:]
    
    content = content.strip()
    
    # Extract image placeholders and descriptions (but DON'T generate yet)
    image_placeholders = extract_image_placeholders(content)
    
    # Map each slide to its image descriptions
    slide_image_descriptions = {}  # slide_index -> list of image descriptions
    
    # Split by slides
    if "Slide" in content:
        slides_content = content.split("Slide")
    else:
        slides_content = content.split("\n## ")

    count = 0
    title_font_color = RGBColor(0, 51, 102) # Dark Blue
    
    # First pass: create all content slides and track which need visual aids
    slides_needing_visuals = []  # List of (slide_index, title, image_descriptions)

    for i, slide_text in enumerate(slides_content):
        if not slide_text.strip(): continue
        
        # Add slide (Title and Content)
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Parse Title vs Body
        lines = slide_text.strip().split('\n')
        # Clean up title: remove "Title:", "#", "*", and whitespace
        title_text = lines[0].replace("Title:", "").replace("title:", "").replace(":", "").replace("#", "").replace("*", "").strip()
        # Remove leading numbering (e.g. "1 Introduction..." -> "Introduction...")
        title_text = title_text.lstrip("0123456789. ")
        body_lines = lines[1:]
        
        # Set Title with Color
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            # Apply styling to title
            if slide.shapes.title.text_frame:
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.color.rgb = title_font_color
                    p.font.bold = True

        # Detect and Process Tables vs Text - collect image descriptions separately
        has_table = False
        table_data = [] # List of rows
        text_content = []
        image_descriptions_for_slide = []  # Image descriptions for this slide
        
        for line in body_lines:
            # Check if this line contains an image placeholder
            is_image_line = False
            for placeholder_info in image_placeholders:
                if placeholder_info['placeholder'] in line:
                    image_descriptions_for_slide.append(placeholder_info['description'])
                    is_image_line = True
                    break
            
            if is_image_line:
                continue  # Skip image placeholder lines
            
            if line.strip().startswith('|') and line.strip().endswith('|'):
                has_table = True
                row = [c.strip() for c in line.strip().strip('|').split('|')]
                if not all(c.replace("-", "") == "" for c in row):
                    table_data.append(row)
            else:
                clean_line = line.strip()
                # Remove image placeholders from text
                for placeholder_info in image_placeholders:
                    clean_line = clean_line.replace(placeholder_info['placeholder'], "")
                
                if clean_line.startswith(("- ", "* ")):
                    clean_line = clean_line[2:]
                clean_line = clean_line.replace("**", "")
                
                if clean_line:
                    lower_line = clean_line.lower()
                    if lower_line.startswith("title:") or lower_line.startswith("slide"):
                         continue
                    text_content.append(clean_line)
        
        # Add Text Content (with limits to prevent overflow)
        MAX_LINES_PER_SLIDE = 6
        MAX_CHARS_PER_LINE = 80
        
        if slide.placeholders[1]:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            # Limit to max lines and truncate long lines
            for line in text_content[:MAX_LINES_PER_SLIDE]:
                truncated_line = line[:MAX_CHARS_PER_LINE] + "..." if len(line) > MAX_CHARS_PER_LINE else line
                p = tf.add_paragraph()
                p.text = truncated_line
                p.level = 0
                p.space_after = Pt(8)
                p.font.size = Pt(16)  # Slightly smaller font for better fit

        # Track slides that need visual aids (limit to 1 image per content slide)
        if image_descriptions_for_slide:
            # Take only the first image description for this slide
            slides_needing_visuals.append((count, title_text, image_descriptions_for_slide[0]))

        # Add Table if present (NO images on content slides)
        if has_table and table_data:
            rows = len(table_data)
            cols = len(table_data[0]) if rows > 0 else 0
            if rows > 0 and cols > 0:
                left = Inches(1)
                top = Inches(3.5)
                width = Inches(8)
                height = Inches(0.8 * rows)
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
                for r_idx, row_data in enumerate(table_data):
                    for c_idx, cell_data in enumerate(row_data):
                        if c_idx < cols:
                            cell = table_shape.cell(r_idx, c_idx)
                            cell.text = cell_data

        count += 1

    # Second pass: Generate images and create dedicated Visual Aid slides
    # Only generate images for a subset to avoid too many API calls (max 3 visual aid slides)
    visuals_to_create = slides_needing_visuals[:3]  # Limit to 3 visual aids per presentation
    
    if visuals_to_create:
        print(f"DEBUG: Creating {len(visuals_to_create)} visual aid slides...")
        
        for slide_idx, slide_title, img_description in visuals_to_create:
            # Generate the image
            print(f"DEBUG: Generating visual aid for: {img_description}")
            img_path = generate_image(img_description, topic_context)
            
            if img_path and os.path.exists(img_path):
                # Create a dedicated visual aid slide (use Title Only layout)
                try:
                    # Use layout 5 (Title Only) or 6 (Blank) for visual aid slides
                    visual_slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
                    visual_slide = prs.slides.add_slide(visual_slide_layout)
                    
                    # Set title for the visual aid slide
                    if visual_slide.shapes.title:
                        visual_slide.shapes.title.text = f"Visual: {slide_title[:40]}"
                        if visual_slide.shapes.title.text_frame:
                            for p in visual_slide.shapes.title.text_frame.paragraphs:
                                p.font.color.rgb = title_font_color
                                p.font.bold = True
                                p.font.size = Pt(28)
                    
                    # Add the image centered on the slide
                    # Standard slide is 10" x 7.5", leave margins
                    img_width = Inches(8)
                    img_height = Inches(5)
                    left = Inches(1)  # Center horizontally: (10 - 8) / 2
                    top = Inches(1.8)  # Below title
                    
                    visual_slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
                    count += 1
                    print(f"DEBUG: Created visual aid slide for '{slide_title}'")
                    
                except Exception as e:
                    print(f"DEBUG: Failed to create visual aid slide: {e}")
            else:
                print(f"DEBUG: Image generation failed for: {img_description}")

    filepath = os.path.join(UPLOAD_DIR, filename)
    prs.save(filepath)
    print(f"DEBUG: PPTX saved to {filepath} ({count} slides)")
    return filepath

def create_pdf_from_text(content: str, filename: str) -> str:
    """Creates a PDF from markdown text."""
    # Remove image placeholders from content (images only for PPTX)
    processed_content = content
    image_placeholders = extract_image_placeholders(content)
    for placeholder_info in image_placeholders:
        processed_content = processed_content.replace(placeholder_info['placeholder'], '')
    
    # Convert MD to HTML
    html = markdown.markdown(processed_content, extensions=['extra', 'codehilite'])
    
    # Enhance HTML
    html_with_styles = f"""
    <html>
    <head>
        <style>
            body {{
                font-family: Arial, sans-serif;
                padding: 20px;
            }}
        </style>
    </head>
    <body>
        {html}
    </body>
    </html>
    """
    
    filepath = os.path.join(UPLOAD_DIR, filename)
    
    with open(filepath, "w+b") as result_file:
        pisa_status = pisa.CreatePDF(html_with_styles, dest=result_file)
        
    return filepath

# --- Global Storage (Singleton Pattern) ---
class IndexManager:
    def __init__(self):
        self.chroma_client = chromadb.PersistentClient(path=PERSIST_DIR)
        
        # Theory Collection
        self.theory_collection = self.chroma_client.get_or_create_collection("theory_collection")
        self.theory_store = ChromaVectorStore(chroma_collection=self.theory_collection)
        self.theory_ctx = StorageContext.from_defaults(vector_store=self.theory_store)
        
        # Lab Collection
        self.lab_collection = self.chroma_client.get_or_create_collection("lab_collection")
        self.lab_store = ChromaVectorStore(chroma_collection=self.lab_collection)
        self.lab_ctx = StorageContext.from_defaults(vector_store=self.lab_store)

        # Indices (Lazy loaded)
        self.theory_index = None
        self.lab_index = None
        self.last_retrieval_error: Optional[str] = None

    def load_indices(self):
        try:
            print(f"[DEBUG] Loading indices from ChromaDB...")
            print(f"[DEBUG] Theory collection count: {self.theory_collection.count()}")
            print(f"[DEBUG] Lab collection count: {self.lab_collection.count()}")
            
            self.theory_index = VectorStoreIndex.from_vector_store(
                self.theory_store, 
                storage_context=self.theory_ctx,
                embed_model=Settings.embed_model
            )
            self.lab_index = VectorStoreIndex.from_vector_store(
                self.lab_store, 
                storage_context=self.lab_ctx,
                embed_model=Settings.embed_model
            )
            print("[DEBUG] Indices loaded from storage successfully!")
        except Exception as e:
            print(f"[ERROR] Indices not found or empty, initializing empty: {e}")
            self.theory_index = VectorStoreIndex.from_documents(
                [], 
                storage_context=self.theory_ctx,
                embed_model=Settings.embed_model
            )
            self.lab_index = VectorStoreIndex.from_documents(
                [], 
                storage_context=self.lab_ctx,
                embed_model=Settings.embed_model
            )

    def get_router_engine(self):
        if not self.theory_index or not self.lab_index:
            self.load_indices()
            
        theory_tool = QueryEngineTool.from_defaults(
            query_engine=self.theory_index.as_query_engine(similarity_top_k=10),
            description="Useful for answering theoretical questions, concepts, lecture notes, lecture slides, and textbook content."
        )
        
        lab_tool = QueryEngineTool.from_defaults(
            query_engine=self.lab_index.as_query_engine(similarity_top_k=10),
            description="Useful for code generation, debugging, lab examples, practical exercises, and programming syntax questions."
        )
        
        return RouterQueryEngine(
            selector=LLMSingleSelector.from_defaults(),
            query_engine_tools=[theory_tool, lab_tool]
        )
    
    def _safe_insert_nodes(self, index, nodes):
        """
        insert_nodes() sometimes raises KeyError when the docstore has been reset
        (e.g. during a rebuild) but nodes still carry a ref_doc_id. In that case
        we embed + write straight to the Chroma vector store so retrieval still
        works. Falling back keeps rebuild_index usable even when llama-index
        docstore bookkeeping is out of sync.
        """
        try:
            index.insert_nodes(nodes)
            return
        except KeyError as e:
            print(f"[WARN] insert_nodes KeyError ({e}); bypassing to vector_store.add")
        except Exception as e:
            print(f"[WARN] insert_nodes failed ({e}); bypassing to vector_store.add")

        try:
            embed_model = Settings.embed_model
            texts = [n.get_content(metadata_mode="none") or "" for n in nodes]
            # Filter out empty/whitespace nodes that break embedding calls
            usable_idx = [i for i, t in enumerate(texts) if t and t.strip()]
            if not usable_idx:
                print("[WARN] No usable node text to embed; skipping vector_store.add")
                return
            usable_nodes = [nodes[i] for i in usable_idx]
            usable_texts = [texts[i] for i in usable_idx]
            embeddings = embed_model.get_text_embedding_batch(usable_texts, show_progress=False)
            final_nodes = []
            for node, emb in zip(usable_nodes, embeddings):
                if isinstance(emb, list) and len(emb) > 0:
                    node.embedding = emb
                    final_nodes.append(node)
            if not final_nodes:
                print("[WARN] Embedding model returned no valid vectors; nothing to add")
                return
            index._vector_store.add(final_nodes)
            print(f"[OK] vector_store.add committed {len(final_nodes)} nodes")
        except Exception as e:
            print(f"[ERROR] vector_store.add fallback also failed: {e}")
            raise

    def ingest_document(self, file_path: str, metadata: dict, file_id: str):
        """
        Ingests a single document into the appropriate index.
        CRITICAL: Ensures doc_id is set to file_id for tracking and future deletion.
        """
        category = metadata.get("category", "theory").lower()
        
        # Sanitize metadata for ChromaDB (no lists/complex types allowed)
        safe_metadata = {}
        for key, value in metadata.items():
            if isinstance(value, (list, dict, tuple)):
                safe_metadata[key] = str(value)
            elif value is None:
                 safe_metadata[key] = ""
            else:
                safe_metadata[key] = value

        # Load document
        documents = SimpleDirectoryReader(input_files=[file_path]).load_data()
        
        for doc in documents:
            # FORCE doc_id to be file_id so we can delete it later
            doc.doc_id = file_id 
            
            # Merge existing doc metadata with our safe metadata
            # We also need to sanitize any new metadata added by the reader (like PPTX slide info)
            combined_metadata = {**doc.metadata, **safe_metadata}
            final_metadata = {}
            for k, v in combined_metadata.items():
                if isinstance(v, (list, dict, tuple)):
                    final_metadata[k] = str(v)
                elif v is None:
                    final_metadata[k] = ""
                else:
                    final_metadata[k] = v
            
            doc.metadata = final_metadata

        nodes = []
        target_index = None
        
        if category == "lab":
            # Lab Index Ingestion with FUNCTION-AWARE Code Splitting
            try:
                ext = os.path.splitext(metadata["filename"])[1]
                language_map = get_language_from_ext(ext)
                print(f"Detected language for {metadata['filename']}: {language_map}")
                
                # Read the actual file content
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        code_content = f.read()
                except:
                    code_content = documents[0].get_content() if documents else ""
                
                # Extract code structure for metadata
                code_structure = extract_code_structure(code_content, language_map)
                print(f"[DEBUG] Extracted code structure: {len(code_structure.get('functions', []))} functions, {len(code_structure.get('classes', []))} classes")
                
                # Try function-aware splitting first (creates one chunk per function)
                function_chunks = split_code_by_functions(code_content, language_map, safe_metadata)
                
                if function_chunks:
                    # Use function-aware chunks - each function is its own document
                    print(f"[DEBUG] Using function-aware splitting: {len(function_chunks)} chunks")
                    
                    # Convert Document objects to nodes
                    from llama_index.core.schema import TextNode
                    nodes = []
                    for i, chunk_doc in enumerate(function_chunks):
                        node = TextNode(
                            text=chunk_doc.text,
                            metadata=chunk_doc.metadata,
                            id_=f"{file_id}_chunk_{i}"
                        )
                        # Add overall code structure for context
                        node.metadata["code_structure"] = json.dumps(code_structure)
                        node.metadata["language"] = language_map
                        nodes.append(node)
                else:
                    # Fallback to standard CodeSplitter for non-Python or if AST fails
                    print(f"[DEBUG] Falling back to CodeSplitter")
                    splitter = CodeSplitter(
                        language=language_map,
                        chunk_lines=60,  # Larger chunks to capture full functions
                        chunk_lines_overlap=20,
                        max_chars=2500
                    )
                    nodes = splitter.get_nodes_from_documents(documents)
                    
                    # Enhance node metadata with code structure
                    for node in nodes:
                        node.metadata["code_structure"] = json.dumps(code_structure)
                        node.metadata["language"] = language_map
                        node.metadata["has_functions"] = str(len(code_structure.get("functions", [])) > 0)
                        node.metadata["has_classes"] = str(len(code_structure.get("classes", [])) > 0)
                        
                        # Add function/class names as searchable text
                        func_names = [f["name"] for f in code_structure.get("functions", [])]
                        class_names = [c["name"] for c in code_structure.get("classes", [])]
                        if func_names:
                            node.metadata["function_names"] = " ".join(func_names)
                        if class_names:
                            node.metadata["class_names"] = " ".join(class_names)
                        
            except Exception as e:
                print(f"Code splitting failed: {e}. Fallback to TokenTextSplitter.")
                import traceback
                traceback.print_exc()
                fallback_splitter = TokenTextSplitter(chunk_size=1024, chunk_overlap=20)
                nodes = fallback_splitter.get_nodes_from_documents(documents)
            
            # Use lab index
            if not self.lab_index: self.load_indices()
            self._safe_insert_nodes(self.lab_index, nodes)
            self.lab_index.storage_context.persist(persist_dir=PERSIST_DIR)
            target_index = self.lab_index
            
        else:
            # Theory Index Ingestion
            # Semantic splitter occasionally blows up on PPTX/PDF chunks where
            # embed model returns an empty vector for near-empty slides. Fall
            # back to a deterministic token splitter so ingest never loses data.
            try:
                splitter = SemanticSplitterNodeParser(
                    buffer_size=1, breakpoint_percentile_threshold=95, embed_model=Settings.embed_model
                )
                nodes = splitter.get_nodes_from_documents(documents)
            except Exception as e:
                print(f"[WARN] SemanticSplitter failed ({e}); falling back to TokenTextSplitter")
                fallback_splitter = TokenTextSplitter(chunk_size=1024, chunk_overlap=80)
                nodes = fallback_splitter.get_nodes_from_documents(documents)
            
            # Add metadata to all nodes (week, topic, course, etc.)
            for node in nodes:
                # Ensure all metadata fields are in node metadata
                for key in ['week', 'topic', 'course', 'category', 'tags', 'filename', 'file_id']:
                    if key in safe_metadata:
                        node.metadata[key] = safe_metadata[key]
            
            if not self.theory_index: self.load_indices()
            self._safe_insert_nodes(self.theory_index, nodes)
            self.theory_index.storage_context.persist(persist_dir=PERSIST_DIR)
            target_index = self.theory_index

        # Ensure all nodes have metadata (for both lab and theory)
        for node in nodes:
            for key in ['week', 'topic', 'course', 'category', 'tags', 'filename', 'file_id']:
                if key in safe_metadata and key not in node.metadata:
                    node.metadata[key] = safe_metadata[key]

        return len(nodes)

    def delete_document(self, file_id: str):
        """
        Removes a document from the indices using its doc_id (which is the file_id).
        """
        if not self.theory_index or not self.lab_index:
            self.load_indices()
            
        # Try deleting from both indices to be safe
        try:
            self.theory_index.delete_ref_doc(file_id, delete_from_docstore=True)
            self.theory_index.storage_context.persist(persist_dir=PERSIST_DIR)
        except Exception:
            pass # Might not exist in theory
            
        try:
            self.lab_index.delete_ref_doc(file_id, delete_from_docstore=True)
            self.lab_index.storage_context.persist(persist_dir=PERSIST_DIR)
        except Exception:
            pass # Might not exist in lab

    def retrieve_context(self, query: str, use_code_aware: bool = True, metadata_filters: Optional[dict] = None) -> List[dict]:
        """
        Retrieves relevant context from both Theory and Lab indices.
        Uses syntax-aware search for code queries when enabled.
        Optionally filters by metadata (week, topic, course, tags, category).
        Returns a list of dicts with 'text' and 'source'.
        """
        if not self.theory_index or not self.lab_index:
            self.load_indices()
            
        self.last_retrieval_error = None
        results = []
        # Retrieval Parameters - Lowered threshold for PPTX slides
        THRESHOLD = 0.35 
        SIMILARITY_TOP_K = 12  # Get more results for metadata filtering
        
        print(f"[DEBUG] Retrieving context for query: {query[:50]}...")
        if metadata_filters:
            print(f"[DEBUG] Metadata filters: {metadata_filters}")
        
        # Load metadata for filtering
        all_metadata = load_metadata() if metadata_filters else {}
        
        # 1. Retrieve from Theory (always semantic)
        try:
            theory_retriever = self.theory_index.as_retriever(similarity_top_k=SIMILARITY_TOP_K)
            theory_nodes = theory_retriever.retrieve(query)
            print(f"[DEBUG] Theory nodes found: {len(theory_nodes)}")
            for node in theory_nodes:
                score = node.score or 0.0
                filename = node.metadata.get("filename", "Unknown Theory Source")
                
                # Apply metadata filtering if filters provided
                if metadata_filters:
                    file_meta = None
                    for file_id, meta in all_metadata.items():
                        if meta.get('filename') == filename:
                            file_meta = meta
                            break
                    
                    if file_meta:
                        # Filter by week (normalized comparison)
                        if 'week' in metadata_filters:
                            file_week = normalize_week(file_meta.get('week', ''))
                            filter_week = normalize_week(metadata_filters['week'])
                            # Extract just the number for comparison
                            file_week_num = re.search(r'(\d+)', file_week)
                            filter_week_num = re.search(r'(\d+)', filter_week)
                            if file_week_num and filter_week_num:
                                if file_week_num.group(1) != filter_week_num.group(1):
                                    continue
                            elif file_week != filter_week and filter_week:
                                continue
                        
                        # Filter by course
                        if 'course' in metadata_filters:
                            file_course = file_meta.get('course', '').lower()
                            filter_course = metadata_filters['course'].lower()
                            if filter_course not in file_course and filter_course:
                                continue
                        
                        # Filter by topic
                        if 'topic' in metadata_filters:
                            file_topic = file_meta.get('topic', '').lower()
                            filter_topic = metadata_filters['topic'].lower()
                            if filter_topic not in file_topic and filter_topic:
                                continue
                        
                        # Filter by category
                        if 'category' in metadata_filters:
                            file_category = file_meta.get('category', '').lower()
                            filter_category = metadata_filters['category'].lower()
                            if file_category != filter_category:
                                continue
                        
                        # Filter by tags
                        if 'tags' in metadata_filters:
                            file_tags = [t.lower() for t in file_meta.get('tags', [])]
                            filter_tags = [t.lower() for t in metadata_filters['tags']]
                            if not any(ft in ' '.join(file_tags) for ft in filter_tags):
                                continue
                
                # Check score if available
                if score >= THRESHOLD:
                    result_item = {
                        "text": node.get_content(),
                        "source": filename,
                        "score": score
                    }
                    # Add metadata to result if available
                    if all_metadata:
                        for file_id, meta in all_metadata.items():
                            if meta.get('filename') == filename:
                                result_item['metadata'] = {
                                    'week': meta.get('week'),
                                    'topic': meta.get('topic'),
                                    'course': meta.get('course'),
                                    'category': meta.get('category'),
                                    'tags': meta.get('tags', [])
                                }
                                break
                    results.append(result_item)
        except Exception as e:
            print(f"[ERROR] Theory retrieval failed: {e}")
            self.last_retrieval_error = str(e)

        # 2. Retrieve from Lab - Use code-aware search if enabled and query is code-related
        try:
            if use_code_aware and is_code_query(query):
                print(f"[DEBUG] Using syntax-aware search for lab materials")
                lab_results = self.retrieve_code_aware(query)
                # Apply metadata filtering to lab results
                if metadata_filters and all_metadata:
                    filtered_lab_results = []
                    for res in lab_results:
                        filename = res['source']
                        file_meta = None
                        for file_id, meta in all_metadata.items():
                            if meta.get('filename') == filename:
                                file_meta = meta
                                break
                        
                        if file_meta:
                            # Apply same metadata filters
                            skip = False
                            if 'week' in metadata_filters:
                                file_week = normalize_week(file_meta.get('week', ''))
                                filter_week = normalize_week(metadata_filters['week'])
                                file_week_num = re.search(r'(\d+)', file_week)
                                filter_week_num = re.search(r'(\d+)', filter_week)
                                if file_week_num and filter_week_num:
                                    if file_week_num.group(1) != filter_week_num.group(1):
                                        skip = True
                                elif file_week != filter_week and filter_week:
                                    skip = True
                            if 'course' in metadata_filters and not skip:
                                if metadata_filters['course'].lower() not in file_meta.get('course', '').lower():
                                    skip = True
                            if 'category' in metadata_filters and not skip:
                                if file_meta.get('category', '').lower() != metadata_filters['category'].lower():
                                    skip = True
                            
                            if not skip:
                                res['metadata'] = {
                                    'week': file_meta.get('week'),
                                    'topic': file_meta.get('topic'),
                                    'course': file_meta.get('course'),
                                    'category': file_meta.get('category'),
                                    'tags': file_meta.get('tags', [])
                                }
                                filtered_lab_results.append(res)
                    results.extend(filtered_lab_results)
                else:
                    results.extend(lab_results)
            else:
                # Standard semantic search
                lab_retriever = self.lab_index.as_retriever(similarity_top_k=SIMILARITY_TOP_K)
                lab_nodes = lab_retriever.retrieve(query)
                print(f"[DEBUG] Lab nodes found: {len(lab_nodes)}")
                for node in lab_nodes:
                    score = node.score or 0.0
                    filename = node.metadata.get("filename", "Unknown Lab Source")
                    
                    # Apply metadata filtering
                    if metadata_filters:
                        file_meta = None
                        for file_id, meta in all_metadata.items():
                            if meta.get('filename') == filename:
                                file_meta = meta
                                break
                        
                        if file_meta:
                            skip = False
                            if 'week' in metadata_filters:
                                file_week = normalize_week(file_meta.get('week', ''))
                                filter_week = normalize_week(metadata_filters['week'])
                                file_week_num = re.search(r'(\d+)', file_week)
                                filter_week_num = re.search(r'(\d+)', filter_week)
                                if file_week_num and filter_week_num:
                                    if file_week_num.group(1) != filter_week_num.group(1):
                                        skip = True
                                elif file_week != filter_week and filter_week:
                                    skip = True
                            if 'course' in metadata_filters and not skip:
                                if metadata_filters['course'].lower() not in file_meta.get('course', '').lower():
                                    skip = True
                            if 'category' in metadata_filters and not skip:
                                if file_meta.get('category', '').lower() != metadata_filters['category'].lower():
                                    skip = True
                            
                            if skip:
                                continue
                    
                    if score >= THRESHOLD:
                        result_item = {
                            "text": node.get_content(),
                            "source": filename,
                            "score": score
                        }
                        # Add metadata
                        if all_metadata:
                            for file_id, meta in all_metadata.items():
                                if meta.get('filename') == filename:
                                    result_item['metadata'] = {
                                        'week': meta.get('week'),
                                        'topic': meta.get('topic'),
                                        'course': meta.get('course'),
                                        'category': meta.get('category'),
                                        'tags': meta.get('tags', [])
                                    }
                                    break
                        results.append(result_item)
        except Exception as e:
            print(f"[ERROR] Lab retrieval failed: {e}")
            self.last_retrieval_error = str(e)
            
        # Sort by score descending
        results.sort(key=lambda x: x["score"], reverse=True)
        return results
    
    def retrieve_code_aware(self, query: str, category: Optional[str] = None) -> List[dict]:
        """
        Syntax-aware search for code materials.
        Combines semantic search with structural matching (function names, classes, patterns).
        """
        if not self.lab_index:
            self.load_indices()
        
        results = []
        THRESHOLD = 0.3  # Lower threshold for code search
        SIMILARITY_TOP_K = 12  # Get more candidates for re-ranking
        
        # Check if this is a code-specific query
        is_code_q = is_code_query(query)
        entities = extract_code_entities(query) if is_code_q else {}
        
        print(f"[DEBUG] Code-aware search - Query: {query[:50]}...")
        print(f"[DEBUG] Is code query: {is_code_q}, Entities: {entities}")
        
        # 1. Standard semantic retrieval
        try:
            lab_retriever = self.lab_index.as_retriever(similarity_top_k=SIMILARITY_TOP_K)
            lab_nodes = lab_retriever.retrieve(query)
            
            for node in lab_nodes:
                base_score = node.score or 0.0
                node_metadata = node.metadata
                
                # 2. Structural matching boost
                structural_boost = 0.0
                match_reasons = []
                
                if is_code_q:
                    # PRIORITY: Check if this chunk IS the function being searched for
                    chunk_name = node_metadata.get("chunk_name", "")
                    chunk_type = node_metadata.get("chunk_type", "")
                    
                    # Heavy boost if chunk_name matches query terms
                    query_terms = query.lower().replace("_", " ").split()
                    if chunk_name and chunk_type == "function":
                        chunk_name_lower = chunk_name.lower().replace("_", " ")
                        # Check for keyword matches in function name
                        for term in ["binary", "search", "linear", "sort", "recursive", 
                                   "fibonacci", "factorial", "merge", "quick", "insertion",
                                   "bubble", "selection", "tower", "hanoi"]:
                            if term in query.lower() and term in chunk_name_lower:
                                structural_boost += 0.3  # Heavy boost for exact function match
                                match_reasons.append(f"exact_function:{chunk_name}")
                                break
                    
                    # Penalize header/footer chunks when looking for implementations
                    if chunk_type in ["header", "footer"]:
                        structural_boost -= 0.2
                    
                    # Match function names from metadata
                    func_names = node_metadata.get("function_names", "")
                    if func_names:
                        for func_name in entities.get("function_names", []):
                            if func_name.lower() in func_names.lower():
                                structural_boost += 0.15
                                match_reasons.append(f"function:{func_name}")
                    
                    # Match class names
                    class_names = node_metadata.get("class_names", "")
                    if class_names:
                        for class_name in entities.get("class_names", []):
                            if class_name in class_names:
                                structural_boost += 0.15
                                match_reasons.append(f"class:{class_name}")
                    
                    # Match patterns
                    code_structure_str = node_metadata.get("code_structure", "{}")
                    try:
                        code_structure = json.loads(code_structure_str)
                        patterns = code_structure.get("patterns", [])
                        for pattern in entities.get("patterns", []):
                            if any(pattern in p for p in patterns):
                                structural_boost += 0.2
                                match_reasons.append(f"pattern:{pattern}")
                    except:
                        pass
                    
                    # Match imports
                    imports_str = " ".join(code_structure.get("imports", [])) if isinstance(code_structure, dict) else ""
                    if imports_str:
                        for imp in entities.get("imports", []):
                            if imp.lower() in imports_str.lower():
                                structural_boost += 0.1
                                match_reasons.append(f"import:{imp}")
                
                # Combine semantic + structural scores
                final_score = base_score + structural_boost
                
                # Only include if meets threshold or has structural matches
                if final_score >= THRESHOLD or match_reasons:
                    results.append({
                        "text": node.get_content(),
                        "source": node_metadata.get("filename", "Unknown Lab Source"),
                        "score": final_score,
                        "base_score": base_score,
                        "structural_boost": structural_boost,
                        "match_reasons": match_reasons,
                        "language": node_metadata.get("language", "unknown"),
                        "chunk_name": node_metadata.get("chunk_name", ""),
                        "chunk_type": node_metadata.get("chunk_type", "")
                    })
                    
        except Exception as e:
            print(f"[ERROR] Code-aware retrieval failed: {e}")
        
        # Sort by final score descending
        results.sort(key=lambda x: x["score"], reverse=True)
        
        # Limit to top results
        return results[:10]
            
    def rebuild_index(self):
        """
        Clears and rebuilds indices from metadata without deleting physical files.
        """
        logger.info("Rebuilding indices...")
        
        # 1. DELETE and RE-CREATE Collections
        try:
            # Delete completely to reset dimension constraints
            try:
                self.chroma_client.delete_collection("theory_collection")
            except Exception:
                pass
            try:
                self.chroma_client.delete_collection("lab_collection")
            except Exception:
                pass
                
            logger.info("Deleted old collections")

            # Re-create fresh
            self.theory_collection = self.chroma_client.get_or_create_collection("theory_collection")
            self.theory_store = ChromaVectorStore(chroma_collection=self.theory_collection)
            self.theory_ctx = StorageContext.from_defaults(vector_store=self.theory_store)
            
            self.lab_collection = self.chroma_client.get_or_create_collection("lab_collection")
            self.lab_store = ChromaVectorStore(chroma_collection=self.lab_collection)
            self.lab_ctx = StorageContext.from_defaults(vector_store=self.lab_store)

            # Re-initialize empty indices
            self.theory_index = VectorStoreIndex.from_documents([], storage_context=self.theory_ctx, embed_model=Settings.embed_model)
            self.lab_index = VectorStoreIndex.from_documents([], storage_context=self.lab_ctx, embed_model=Settings.embed_model)
            
        except Exception as e:
            logger.error(f"Error resetting indices: {e}")
            raise e

        # 2. Re-ingest all files
        all_metadata = load_metadata()
        count = 0
        errors = []
        details = []
        total = len(all_metadata)
        
        for file_id, meta in all_metadata.items():
            try:
                filename = meta.get("filename")
                file_path = os.path.join(UPLOAD_DIR, filename)
                
                if not os.path.exists(file_path):
                    msg = f"File not found: {filename}"
                    logger.warning(msg)
                    errors.append(msg)
                    continue
                    
                self.ingest_document(file_path, meta, file_id)
                count += 1
                msg = f"Re-ingested {count}/{total}: {filename}"
                logger.info(msg)
                details.append(msg)
                
            except Exception as e:
                msg = f"Failed to re-ingest {file_id}: {str(e)}"
                logger.error(msg)
                errors.append(msg)
                
        return {"count": count, "errors": errors, "details": details}

# Initialize Manager
index_manager = IndexManager()

def generate_material_core(prompt: str, category: str, type: str):
    """
    Core logic for generating study materials.
    Returns dict with content, downloadUrl, filename.
    """
    # 1. Retrieve Internal Context (RAG)
    print(f"DEBUG: Retrieving context for generation query: {prompt}")
    context_data = index_manager.retrieve_context(prompt)
    context_str = ""
    if context_data:
        context_str = "## RELEVANT COURSE MATERIALS (Use these as primary source):\n"
        for item in context_data[:5]: # Top 5 sources
            context_str += f"---\nSOURCE: {item['source']}\nCONTENT: {item['text'][:1500]}\n---\n"
    
    # Construct a tailored system prompt
    base_instruction = """
    You are an expert educational content generator for a university-level platform.
    Output MUST be in Markdown format.
    """
    
    specific_instruction = ""
    output_ext = ".md"
    actual_type = type
    
    print(f"DEBUG: generate_material_core called with type='{type}', category='{category}'")
    
    # Handle visual/image generation separately (no LLM content generation needed)
    if type.lower().strip() in ['visual', 'image']:
        print(f"DEBUG: Generating standalone visual for: {prompt}")
        img_path = generate_image(f"Educational diagram or visual aid about: {prompt}", prompt)
        
        if img_path and os.path.exists(img_path):
            img_filename = os.path.basename(img_path)
            return {
                "content": "[IMAGE_GENERATED]",
                "downloadUrl": f"/static/materials/{img_filename}",
                "filename": img_filename,
                "isVisual": True
            }
        else:
            return {
                "content": "Failed to generate visual. Please try again with a more specific description.",
                "downloadUrl": None,
                "filename": None,
                "isVisual": True
            }
    
    if type == 'code':
        if not is_code_request(prompt):
            return {
                "content": "Not a valid question for code generation. Please ask for code specifically.",
                "downloadUrl": None,
                "filename": None
            }
        else:
            lang_name, lang_ext = detect_language_from_prompt(prompt)
            specific_instruction = f"""
            CRITICAL: Return ONLY valid, runnable {lang_name.upper()} code. ABSOLUTELY NO conversational text.
            - Do NOT write the language name at the start.
            - Do NOT say "Here is the code".
            - Do NOT use markdown code blocks (```). Return raw code only.
            - Include comments WITHIN the code to explain logic.
            """
            output_ext = lang_ext
            actual_type = 'code'
    
    elif type == 'reading':
        specific_instruction = """
        Generate CODE-CENTRIC learning material as a structured document:
        - Format with clear H1, H2 headers.
        - Include SHORT code snippets (5-15 lines) to illustrate concepts.
        - Keep explanations concise but thorough.
        """
        output_ext = ".pdf"
        actual_type = 'reading'
    
    elif type == 'slides':
        specific_instruction = """
        PROMPT FOR SLIDES - STRICT FORMATTING RULES:
        1. Use "Slide [N]: [Title]" as the header for each slide.
        2. MAXIMUM 5-6 bullet points per slide. NO MORE.
        3. Each bullet point must be SHORT (max 12 words). Be concise!
        4. Do NOT use nested lists or sub-bullets.
        5. MATH: Use UNICODE symbols (→, ∑, ≤), NOT LaTeX.
        6. Do NOT include conversational text.
        7. Split content across MORE slides if needed - brevity is key!
        8. VISUAL AIDS: Include 2-3 [IMAGE: description] placeholders total.
        """

        output_ext = ".pptx"
        actual_type = 'slides'
    
    elif type == 'pdf' or category == 'theory':
            specific_instruction = """
            Format as a structured document with clear spacing.
            - Use H1 (#) ONLY for the main document title.
            - Use H2 (##) for major sections.
            - Add empty lines between paragraphs.
            """
            output_ext = ".pdf" if type == 'pdf' else ".md"
            actual_type = type

    full_prompt = f"""{base_instruction}
{specific_instruction}

{context_str}

Context: {category} ({actual_type})
User Prompt: {prompt}
"""
    
    CHAT_MODELS = [
        # 'gemini-2.0-flash',
        # 'gemini-2.0-flash-lite',
        'gemini-2.5-flash',
        'gemini-3-flash-preview',
    ]

    
    content = ""
    last_error = None
    
    for model_name in CHAT_MODELS:
        try:
            # Use the Official Google GenAI SDK
            client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
            
            # Tools: Google Search for external context
            tools = [types.Tool(google_search=types.GoogleSearch())]
            
            print(f"DEBUG: Generating content using model={model_name}...")
            response = client.models.generate_content(
                model=model_name,
                contents=full_prompt,
                config=types.GenerateContentConfig(
                    tools=tools,
                    temperature=0.4 # Lower temperature for more grounded content
                )
            )
            content = response.text
            
            # Check for grounding
            if response.candidates and response.candidates[0].grounding_metadata:
                print("DEBUG: Grounding metadata received.")
            
            if content:
                break # Success!

        except Exception as e:
            error_str = str(e)
            print(f"DEBUG: Model {model_name} failed: {error_str[:100]}")
            last_error = e
            # Retry with next model for server issues OR if model doesn't support tools
            if "503" in error_str or "429" in error_str or "overloaded" in error_str.lower() or "unsupported" in error_str.lower():
                continue
            else:
                break # Non-retryable error

    if not content:
        if last_error:
            # Final fallback to a very basic call if all else fails
            try:
                print("DEBUG: Final fallback attempt...")
                client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
                response = client.models.generate_content(
                    model='gemini-2.0-flash-lite', 
                    contents=full_prompt
                )
                content = response.text
            except Exception as e2:
                 raise HTTPException(status_code=500, detail=f"Generation failed completely: {e2}")
        else:
             raise HTTPException(status_code=500, detail="Content generation returned empty response.")

    # Cleanup artifacts
    content = re.sub(r'```\w*\n?', '', content).replace("```", "").strip()

    # File Generation Logic
    filename = f"{actual_type}_{uuid.uuid4().hex[:8]}{output_ext}"
    download_url = None
    
    if actual_type == 'slides':
        create_pptx_from_text(content, filename, topic_context=prompt)
        download_url = f"/static/materials/{filename}"
        content = f"[System] Presentation generated. Download below.\n\nPreview:\n{content[:500]}..."
        
    elif actual_type == 'pdf':
        create_pdf_from_text(content, filename)
        download_url = f"/static/materials/{filename}"
    
    elif actual_type == 'notes':
            filepath = os.path.join(UPLOAD_DIR, filename)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            download_url = f"/static/materials/{filename}"
    
    elif actual_type == 'reading':
            create_pdf_from_text(content, filename)
            download_url = f"/static/materials/{filename}"
    
    if actual_type == 'code':
            lang_name_detected, _ = detect_language_from_prompt(prompt) 
            content = clean_code_content(content, lang_name_detected)
            
            # --- validation start ---
            validation = validate_code_with_fallback(content, lang_name_detected)
            if not validation['valid']:
                print(f"DEBUG: Code validation failed. Error: {validation['error'][:100] if validation.get('error') else 'Unknown error'}...")
                # SELF-CORRECTION: Retry once
                print("DEBUG: Attempting Self-Correction...")
                
                correction_prompt = f"""
The previous code you generated for '{lang_name_detected}' failed to run/compile.
Error Message:
{validation['error']}

Fix the code and return ONLY the valid, runnable code.
"""
                
                retry_prompt = f"{full_prompt}\n\nPREVIOUS_OUTPUT:\n{content}\n\nERROR:\n{correction_prompt}"
                
                try:
                    client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
                    response = client.models.generate_content(
                        model=chosen_model,
                        contents=retry_prompt
                    )
                    fixed_content = response.text
                    content = clean_code_content(fixed_content, lang_name_detected)
                    content = re.sub(r'```\w*\n?', '', content).replace("```", "").strip()
                    print("DEBUG: Self-Correction successful (new content generated).")
                    
                    # Re-validate the fixed code
                    re_validation = validate_code_with_fallback(content, lang_name_detected)
                    if not re_validation['valid']:
                        error_comment = f"\n\n/* \n[System] Code Verification Failed even after correction.\nError: {re_validation['error']}\n*/"
                        if lang_name_detected == 'python':
                            error_comment = f"\n\n'''\n[System] Code Verification Failed even after correction.\nError: {re_validation['error']}\n'''"
                        content += error_comment
                        
                except Exception as e:
                    print(f"DEBUG: Self-Correction failed: {e}")
                    # Append original error
                    error_comment = f"\n\n/* \n[System] Code Verification Failed.\nError: {validation['error']}\n*/"
                    if lang_name_detected == 'python':
                        error_comment = f"\n\n'''\n[System] Code Verification Failed.\nError: {validation['error']}\n'''"
                    content += error_comment

            # --- validation end ---
            
            filepath = os.path.join(UPLOAD_DIR, filename)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            download_url = f"/static/materials/{filename}"

    return {
        "content": content,
        "downloadUrl": download_url,
        "filename": filename
    }

def generate_study_material_tool(topic: str, format: str):
    """
    Generates downloadable study materials and returns a download link.
    Args:
        topic: The detailed topic or prompt for the content.
        format: One of 'slides', 'pdf', 'code', 'reading' (lab guide), 'notes', 'visual' (image/diagram).
    """
    # Default logic for category
    category = "lab" if format in ['code', 'reading'] else "theory"
    result = generate_material_core(topic, category, format)
    
    # Return explicit markdown for the model to use
    return f"Material generated successfully. [Download {format.upper()} File]({result['downloadUrl']}) "

# --- App Initialization ---
@asynccontextmanager
async def lifespan(app: FastAPI):
    index_manager.load_indices()
    yield

app = FastAPI(title="StudySmart Backend", lifespan=lifespan)

# CORS Configuration - Updated to be more permissive for development
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # Allow all origins for dev simplicity
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Include forum + analytics routers
from forum import router as forum_router
app.include_router(forum_router)
app.include_router(analytics_router)

# Serve uploaded files statically
os.makedirs(UPLOAD_DIR, exist_ok=True)
app.mount("/static/materials", StaticFiles(directory=UPLOAD_DIR), name="materials")

# --- Data Models ---
class ChatRequest(BaseModel):
    message: str
    history: Optional[List[dict]] = []
    allow_web_search: bool = False

class GenerateRequest(BaseModel):
    prompt: str
    category: str
    type: str

class ValidateRequest(BaseModel):
    code: str
    language: str

# --- Endpoints ---

@app.post("/api/upload")
async def upload_file(
    file: UploadFile = File(...), 
    category: str = Form(...),
    course: str = Form(None),
    topic: str = Form(None),
    week: str = Form(None),
    tags: str = Form(None),
    user=Depends(require_teacher),
):
    try:
        # 1. Save File Locally
        file_path = os.path.join(UPLOAD_DIR, file.filename)
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        stat = os.stat(file_path)
        file_id = file.filename # Using filename as durable ID
        
        # 2. Prepare Metadata
        metadata_entry = {
            "id": file_id,
            "filename": file.filename,
            "title": os.path.splitext(file.filename)[0].replace("_", " ").replace("-", " "),
            "category": category,
            "course": course or "General",
            "topic": topic or "General",
            "week": week or "",
            "tags": [t.strip() for t in (tags or "").split(",") if t.strip()],
            "type": file.filename.split('.')[-1].lower(),
            "size": f"{stat.st_size / (1024 * 1024):.2f} MB",
            "size_bytes": stat.st_size,
            "date": datetime.now().strftime("%Y-%m-%d"),
            "isValidated": True,
            "uploadedBy": user.get("uid"),
            "uploadedAt": datetime.utcnow().isoformat(),
        }
        
        # Clean metadata specifically for embedding (flatten lists)
        embed_metadata = {
            "filename": file.filename,
            "file_id": file_id,
            "category": category,
            "course": course or "General",
            "topic": topic or "General",
            "week": week or "",
            "tags": tags or ""
        }

        # 3. Store Metadata
        repo.upsert_material(file_id, metadata_entry)

        # 4. Ingest into Vector Store (removing old version if exists first to be safe)
        index_manager.delete_document(file_id)
        nodes_count = index_manager.ingest_document(file_path, embed_metadata, file_id)

        # 5. Log analytics event
        repo.log_event(
            event_type="upload",
            uid=user.get("uid"),
            role=user.get("role"),
            metadata={"fileId": file_id, "filename": file.filename, "target": metadata_entry["title"]},
        )

        return {
            "status": "success", 
            "id": file_id,
            "filename": file.filename,
            "nodes_ingested": nodes_count
        }

    except Exception as e:
        print(f"Upload failed: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/validate")
async def validate_code_endpoint(request: ValidateRequest):
    """Validates code syntax and compilation."""
    print(f"DEBUG: Validating code manually (Language: {request.language})...")
    result = validate_code_with_fallback(request.code, request.language)
    return result

@app.get("/api/courses")
async def list_courses():
    """Returns unique course names from all uploaded materials."""
    try:
        all_metadata = load_metadata()
        courses = set()
        for material in all_metadata.values():
            course_name = material.get("course", "").strip()
            if course_name and course_name != "General":
                courses.add(course_name)
        return {"courses": sorted(list(courses))}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/materials")
async def list_materials():
    try:
        all_metadata = load_metadata()
        materials = list(all_metadata.values())
        # Sort by date descending
        materials.sort(key=lambda x: x.get("date", ""), reverse=True)
        return materials
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/materials/{file_id}")
async def delete_material(file_id: str, user=Depends(require_teacher)):
    try:
        material = repo.get_material(file_id)
        if not material:
            raise HTTPException(status_code=404, detail="Material not found")

        # 1. Delete from Vector Index
        index_manager.delete_document(file_id)

        # 2. Delete File
        file_path = os.path.join(UPLOAD_DIR, material.get("filename", file_id))
        if os.path.exists(file_path):
            os.remove(file_path)

        # 3. Remove Metadata from Firestore
        repo.delete_material(file_id)

        return {"status": "success", "deleted": file_id}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/materials/{file_id}")
async def update_material(file_id: str, updates: dict, user=Depends(require_teacher)):
    """
    Updates metadata AND re-indexes the document so RAG reflects changes.
    """
    try:
        current_meta = repo.get_material(file_id)
        if not current_meta:
            raise HTTPException(status_code=404, detail="Material not found")

        # Apply updates to allowed fields
        allowed_fields = ["title", "course", "topic", "week", "tags", "category"]
        needs_reindex = False
        patch: dict = {}

        for field in allowed_fields:
            if field in updates and updates[field] != current_meta.get(field):
                current_meta[field] = updates[field]
                patch[field] = updates[field]
                needs_reindex = True

        # Save metadata to Firestore
        if patch:
            repo.update_material(file_id, patch)

        # Re-index if critical metadata changed (CRITICAL FIX)
        if needs_reindex:
            file_path = os.path.join(UPLOAD_DIR, current_meta["filename"])
            if os.path.exists(file_path):
                # Flatten tags for embedding
                tags_str = ",".join(current_meta["tags"]) if isinstance(current_meta["tags"], list) else current_meta["tags"]
                
                embed_metadata = {
                    "filename": current_meta["filename"],
                    "file_id": file_id,
                    "category": current_meta["category"],
                    "course": current_meta["course"],
                    "topic": current_meta["topic"],
                    "week": current_meta["week"],
                    "tags": tags_str
                }
                
                # Delete old vectors and re-ingest
                index_manager.delete_document(file_id)
                index_manager.ingest_document(file_path, embed_metadata, file_id)
        
        return current_meta
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/admin/reindex")
async def reindex_materials(user=Depends(require_teacher)):
    """Triggers a full re-index of all materials."""
    try:
        result = index_manager.rebuild_index()
        return {"status": "success", "result": result}
    except Exception as e:
        logger.error(f"Re-indexing failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/chat")
async def chat(request: ChatRequest, user=Depends(optional_current_user)):
    try:
        if user:
            repo.log_event(
                event_type="chat",
                uid=user.get("uid"),
                role=user.get("role"),
                metadata={"target": request.message[:120]},
            )
        # 0. Smart Greeting / Conversational Check
        greetings = ["hi", "hello", "hey", "greetings", "good morning", "good afternoon", "good evening", "yo", "what's up", "sup"]
        if request.message.lower().strip() in greetings:
            # Load metadata to show available materials
            all_metadata = load_metadata()
            materials_summary = ""
            if all_metadata:
                weeks = sorted(set(str(m.get('week', '')).strip() for m in all_metadata.values() if m.get('week')))
                courses = sorted(set(m.get('course', '') for m in all_metadata.values() if m.get('course')))
                if weeks:
                    materials_summary = f"\n\n📚 **Available Materials:**\n"
                    materials_summary += f"• Weeks: {', '.join(weeks[:10])}{'...' if len(weeks) > 10 else ''}\n"
                    if courses:
                        materials_summary += f"• Courses: {', '.join(courses[:5])}{'...' if len(courses) > 5 else ''}\n"
                    materials_summary += f"• Total files: {len(all_metadata)}\n"
                    materials_summary += "\nYou can ask about specific weeks, topics, or courses!"
            
            return {
                "response": "Hello! I'm StudySmart, your AI Study Assistant. I can help you:\n\n"
                           "• **Search course materials** - Ask about any topic\n"
                           "• **Summarize documents** - Say 'summarize [topic]'\n"
                           "• **Explain concepts** - Ask 'explain [concept]'\n"
                           "• **Answer questions** - Based on your uploaded materials\n"
                           "• **Find materials by week/topic/course** - e.g., 'what are the week 7 contents'\n\n"
                           "What would you like to learn about today?" + materials_summary, 
                "sources": []
            }
        
        # 0.5. Extract metadata filters from query
        metadata_filters = extract_metadata_filters(request.message)
        print(f"[DEBUG] Extracted metadata filters: {metadata_filters}")
        
        # Check if this is a metadata-only query (e.g., "what are the week 7 contents")
        is_metadata_query = bool(metadata_filters) and any(key in metadata_filters for key in ['week', 'topic', 'course'])
        query_lower = request.message.lower()
        is_listing_query = any(phrase in query_lower for phrase in [
            'what are', 'show me', 'list', 'what contents', 'what materials', 
            'what files', 'what documents', 'contents of', 'materials for'
        ])
        
        # Initialize sources tracking
        sources = []
        seen_sources = set()
        
        # 1. Retrieve Context from Vector Store with metadata filtering
        results = index_manager.retrieve_context(request.message, use_code_aware=True, metadata_filters=metadata_filters if metadata_filters else None)
        
        # 1.5. If it's a metadata listing query, always check metadata directly to ensure completeness
        metadata_context = ""
        matching_materials = []
        if is_metadata_query and is_listing_query:
            all_metadata = load_metadata()
            for file_id, meta in all_metadata.items():
                match = True
                if 'week' in metadata_filters:
                    file_week = normalize_week(meta.get('week', ''))
                    filter_week = normalize_week(metadata_filters['week'])
                    file_week_num = re.search(r'(\d+)', file_week)
                    filter_week_num = re.search(r'(\d+)', filter_week)
                    if file_week_num and filter_week_num:
                        if file_week_num.group(1) != filter_week_num.group(1):
                            match = False
                    elif file_week != filter_week and filter_week:
                        match = False
                if 'course' in metadata_filters and match:
                    if metadata_filters['course'].lower() not in meta.get('course', '').lower():
                        match = False
                if 'topic' in metadata_filters and match:
                    if metadata_filters['topic'].lower() not in meta.get('topic', '').lower():
                        match = False
                if 'category' in metadata_filters and match:
                    if meta.get('category', '').lower() != metadata_filters['category'].lower():
                        match = False
                
                if match:
                    matching_materials.append(meta)
            
            # Create context from metadata - always show all matching materials for listing queries
            if matching_materials:
                print(f"[DEBUG] Found {len(matching_materials)} materials by metadata for query: {metadata_filters}")
                metadata_context = f"MATERIALS MATCHING YOUR QUERY ({metadata_filters}):\n\n"
                # Show all matching materials
                for mat in matching_materials:
                    filename = mat.get('filename', '')
                    metadata_context += f"- **{mat.get('title', filename)}**\n"
                    metadata_context += f"  - Week: {mat.get('week', 'N/A')}\n"
                    metadata_context += f"  - Topic: {mat.get('topic', 'N/A')}\n"
                    metadata_context += f"  - Course: {mat.get('course', 'N/A')}\n"
                    metadata_context += f"  - Category: {mat.get('category', 'N/A')}\n"
                    if mat.get('tags'):
                        metadata_context += f"  - Tags: {', '.join(mat.get('tags', []))}\n"
                    metadata_context += f"  - Type: {mat.get('type', 'N/A')}\n"
                    metadata_context += f"  - Filename: {filename}\n\n"
                    
                    # Add to sources if not already added
                    if filename and filename not in seen_sources:
                        sources.append({
                            "title": mat.get('title', filename),
                            "type": "file",
                            "url": f"http://localhost:8000/static/materials/{filename}"
                        })
                        seen_sources.add(filename)

        # 1.6. Load all metadata to provide context about available materials
        all_metadata = load_metadata()
        metadata_summary = ""
        if all_metadata:
            # Group by week, course, topic for better context
            weeks_dict = {}
            courses_dict = {}
            topics_dict = {}
            
            for file_id, meta in all_metadata.items():
                week = meta.get('week', '').strip()
                course = meta.get('course', '').strip()
                topic = meta.get('topic', '').strip()
                
                if week:
                    if week not in weeks_dict:
                        weeks_dict[week] = []
                    weeks_dict[week].append(meta.get('title', meta.get('filename', '')))
                
                if course:
                    if course not in courses_dict:
                        courses_dict[course] = []
                    courses_dict[course].append(meta.get('title', meta.get('filename', '')))
                
                if topic:
                    if topic not in topics_dict:
                        topics_dict[topic] = []
                    topics_dict[topic].append(meta.get('title', meta.get('filename', '')))
            
            metadata_summary = "\n\n## AVAILABLE MATERIALS METADATA:\n"
            if weeks_dict:
                metadata_summary += f"**Weeks Available:** {', '.join(sorted(weeks_dict.keys(), key=lambda x: int(x) if x.isdigit() else 999))}\n"
            if courses_dict:
                metadata_summary += f"**Courses:** {', '.join(sorted(courses_dict.keys())[:10])}\n"
            if topics_dict:
                metadata_summary += f"**Topics:** {', '.join(sorted(topics_dict.keys())[:15])}\n"
        
        # Build context string
        context_str = ""
        if metadata_context:
            context_str = metadata_context
            # Add metadata summary for additional context
            if metadata_summary:
                context_str += "\n" + metadata_summary
        
        if results:
            if context_str:
                context_str += "\n\nRELEVANT COURSE MATERIALS CONTENT:\n"
            else:
                context_str = "RELEVANT COURSE MATERIALS:\n"
            
            for res in results:
                # Include more content for better context
                content_preview = res['text'][:2000] if len(res['text']) > 2000 else res['text']
                metadata_info = ""
                if 'metadata' in res:
                    meta = res['metadata']
                    metadata_info = f"METADATA: Week={meta.get('week', 'N/A')}, Topic={meta.get('topic', 'N/A')}, Course={meta.get('course', 'N/A')}, Category={meta.get('category', 'N/A')}"
                
                context_str += f"---\nSOURCE: {res['source']}\n{metadata_info}\nSCORE: {res.get('score', 'N/A')}\nCONTENT:\n{content_preview}\n---\n\n"
                
                filename = res['source']
                if filename not in seen_sources:
                    sources.append({
                        "title": filename,
                        "type": "file",
                        "url": f"http://localhost:8000/static/materials/{filename}"
                    })
                    seen_sources.add(filename)
        else:
            # If no results but metadata filters were used, provide helpful context
            if metadata_filters and not metadata_context:
                context_str = f"No materials found matching the filters: {metadata_filters}.\n" + metadata_summary
            elif not context_str:
                context_str = "No directly relevant materials found for this query.\n" + metadata_summary
            elif metadata_context and metadata_summary:
                # Ensure metadata summary is included even when we have metadata_context
                context_str += "\n" + metadata_summary
        
        # 2. Setup Gemini Client
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
        
        # 3. Configure Tools & Instructions based on settings
        # IMPORTANT: You CANNOT mix google_search with custom function tools - they are mutually exclusive!
        web_search_instruction = ""
        
        if request.allow_web_search:
            web_search_instruction = (
                "   b) If course materials are insufficient or unavailable, use Google Search to find accurate information.\n"
                "   c) When using web search, clearly cite the source with a link.\n"
                "   d) Note: File generation is disabled when web search is enabled.\n"
            )
            # Use ONLY Google Search tool (cannot mix with custom functions)
            tools = [types.Tool(google_search=types.GoogleSearch())]
        else:
            web_search_instruction = (
                "   b) If course materials are insufficient, clearly state that and suggest enabling web search or provide general guidance.\n"
            )
            # Use custom function tool for generating materials
            tools = [generate_study_material_tool]
        
        # 4. Enhanced System Prompt for Part 5 Requirements (Metadata-Aware)
        system_instruction = (
            "You are 'StudySmart', an intelligent AI Study Assistant for university courses.\n\n"
            
            "## YOUR CORE CAPABILITIES:\n"
            "1. **SEARCH & RETRIEVE**: Find relevant information from uploaded course materials (slides, PDFs, code, notes)\n"
            "2. **METADATA QUERIES**: Answer questions about materials by week, topic, course, or tags (e.g., 'what are the week 7 contents', 'show me materials on neural networks')\n"
            "3. **SUMMARIZE**: Create concise summaries of topics, documents, or concepts\n"
            "4. **EXPLAIN**: Provide clear, educational explanations of complex topics\n"
            "5. **ANSWER QUESTIONS**: Give accurate, grounded responses based on course content\n"
            "6. **GENERATE FILES**: If user asks to create/generate slides, notes, reports, PDFs, or code files, IMMEDIATELY CALL the `generate_study_material` tool with the topic extracted from their request. DO NOT ask for confirmation - just generate it directly.\n"
            "7. **FOLLOW-UP**: Handle contextual follow-up questions using conversation history\n\n"
            
            "## METADATA AWARENESS:\n"
            "The system has access to materials with metadata including:\n"
            "- **Week**: Materials organized by week (e.g., Week 1, Week 7)\n"
            "- **Topic**: Subject matter or topic tags\n"
            "- **Course**: Course name or identifier\n"
            "- **Category**: Theory or Lab\n"
            "- **Tags**: Additional descriptive tags\n\n"
            
            "When users ask about specific weeks, topics, or courses, you should:\n"
            "- Reference the metadata information provided\n"
            "- List materials matching their query\n"
            "- Provide summaries or details about those materials\n"
            "- If materials exist but weren't retrieved, mention them from the metadata summary\n\n"
            
            "## INSTRUCTIONS:\n"
            "1. **GROUNDING FIRST**: Always prioritize information from 'RELEVANT COURSE MATERIALS' when available.\n"
            "2. **METADATA QUERIES**: If user asks 'what are the week X contents' or similar metadata queries:\n"
            "   - Check the AVAILABLE MATERIALS METADATA section\n"
            "   - List all materials for that week/topic/course\n"
            "   - Provide summaries or details about those materials\n"
            "   - If materials exist in metadata but weren't retrieved, still mention them\n"
            "3. **CITATIONS**: When using course materials, cite the specific filename and metadata (e.g., 'According to `lecture5.pdf` (Week 7, Topic: Neural Networks)...')\n"
            f"{web_search_instruction}"
            "4. **FORMATTING**: Use markdown for clarity - headers, bullet points, code blocks where appropriate.\n"
            "5. **CODE**: When showing code, ensure it is syntactically correct and properly formatted.\n"
            "6. **CONTEXT-AWARE**: Use the conversation history to understand follow-up questions.\n"
            "7. **EDUCATIONAL TONE**: Be helpful, encouraging, and educational.\n\n"
            
            "## AVAILABLE COURSE MATERIALS:\n"
            f"{context_str if context_str else 'No directly relevant materials found for this query.'}\n"
        )

        # 5. Build Conversation History for Multi-Turn Context
        conversation_contents = []
        
        # Add previous messages from history
        if request.history:
            for msg in request.history[-10:]:  # Limit to last 10 messages for context window
                role = msg.get("role", "user")
                content = msg.get("content", "")
                if role == "assistant":
                    conversation_contents.append(
                        types.Content(role="model", parts=[types.Part.from_text(text=content)])
                    )
                else:
                    conversation_contents.append(
                        types.Content(role="user", parts=[types.Part.from_text(text=content)])
                    )
        
        # Add current message
        conversation_contents.append(
            types.Content(role="user", parts=[types.Part.from_text(text=request.message)])
        )

        # Model fallback list - ONLY models that support function calling/tools
        CHAT_MODELS = [
            'gemini-2.0-flash',  # Best for function calling
            'gemini-2.0-flash-lite',  # Fallback
        ]
        
        last_error = None
        for model_name in CHAT_MODELS:
            try:
                print(f"[DEBUG] Trying model: {model_name}")
                # 6. Generate Response with Full Context
                response = client.models.generate_content(
                    model=model_name,
                    contents=conversation_contents,
                    config=types.GenerateContentConfig(
                        system_instruction=system_instruction,
                        tools=tools,
                        temperature=0.7 
                    )
                )
                
                response_text = response.text if response.text else "I apologize, but I couldn't generate a response. Please try rephrasing your question."
                
                # 7. Extract Web Sources if available
                if response.candidates and response.candidates[0].grounding_metadata and response.candidates[0].grounding_metadata.grounding_chunks:
                    for chunk in response.candidates[0].grounding_metadata.grounding_chunks:
                        if chunk.web and chunk.web.title and chunk.web.uri:
                            if chunk.web.uri not in seen_sources:
                                sources.append({
                                    "title": chunk.web.title,
                                    "type": "web",
                                    "url": chunk.web.uri
                                })
                                seen_sources.add(chunk.web.uri)
                
                print(f"[DEBUG] Successfully used model: {model_name}")
                return {"response": response_text, "sources": sources}

            except Exception as e:
                error_str = str(e)
                print(f"[WARN] Model {model_name} failed: {error_str[:100]}")
                last_error = e
                
                # Retry with next model for server issues OR if model doesn't support tools
                if "503" in error_str or "429" in error_str or "overloaded" in error_str.lower() or "unsupported" in error_str.lower():
                    continue  # Try next model
                else:
                    # For other errors, don't retry - raise immediately
                    raise e
        
        # All models failed - return graceful error
        if last_error:
            if "429" in str(last_error) or "503" in str(last_error):
                return {
                    "response": "⏳ All AI models are currently experiencing high traffic. Please wait 30-60 seconds and try again.\n\n"
                               "**Tip:** If this persists, consider linking a billing account to your Google Cloud project for higher quotas.",
                    "sources": []
                }
            raise last_error

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

# --- PART 2: Intelligent Search Engine ---
class SearchRequest(BaseModel):
    query: str
    category: Optional[str] = None  # "theory", "lab", or None for all
    course: Optional[str] = None    # Filter by specific course
    limit: int = 10

@app.post("/api/search")
async def search_materials(request: SearchRequest):
    """
    Intelligent search endpoint for course materials (Part 2).
    - Semantic search beyond keyword matching ✓
    - RAG-based retrieval ✓
    - Syntax-aware/structure-aware search for lab/code materials (BONUS) ✓
    Returns relevant documents, excerpts, or code snippets.
    """
    try:
        # Determine if we should use code-aware search
        use_code_aware = request.category is None or request.category.lower() == "lab"
        
        # Perform intelligent retrieval (semantic + syntax-aware for code)
        results = index_manager.retrieve_context(request.query, use_code_aware=use_code_aware)

        # Surface real backend problems instead of masking them as "no results".
        theory_count = 0
        lab_count = 0
        try:
            theory_count = index_manager.theory_collection.count()
            lab_count = index_manager.lab_collection.count()
        except Exception:
            pass

        if not results:
            last_error = index_manager.last_retrieval_error or ""
            if "expecting embedding with dimension" in last_error.lower():
                raise HTTPException(
                    status_code=503,
                    detail=(
                        "Search index is out of date (embedding dimension mismatch). "
                        "Ask an admin to run POST /api/admin/reindex to rebuild vector indices."
                    ),
                )
            if last_error:
                # Any other retrieval-time failure (API quota, network, etc.)
                raise HTTPException(
                    status_code=503,
                    detail=f"Search backend error: {last_error[:300]}",
                )
            if theory_count == 0 and lab_count == 0:
                raise HTTPException(
                    status_code=503,
                    detail=(
                        "Search index is empty — no materials have been embedded yet. "
                        "Upload materials or run POST /api/admin/reindex to rebuild."
                    ),
                )
            return {
                "query": request.query,
                "results": [],
                "total": 0,
                "message": "No matching materials found. Try different keywords or upload more content.",
                "search_type": "semantic" + (" + syntax-aware" if use_code_aware and is_code_query(request.query) else "")
            }
        
        # Filter by category if specified
        if request.category:
            # Load metadata to get categories
            all_metadata = load_metadata()
            filtered_results = []
            for res in results:
                filename = res['source']
                # Find metadata for this file
                for file_id, meta in all_metadata.items():
                    if meta.get('filename') == filename:
                        if meta.get('category', '').lower() == request.category.lower():
                            filtered_results.append(res)
                        break
            results = filtered_results
        
        # Filter by course if specified
        if request.course:
            all_metadata = load_metadata()
            filtered_results = []
            for res in results:
                filename = res['source']
                for file_id, meta in all_metadata.items():
                    if meta.get('filename') == filename:
                        if meta.get('course', '').lower() == request.course.lower():
                            filtered_results.append(res)
                        break
            results = filtered_results
        
        # Limit results
        results = results[:request.limit]
        
        # Format response with enhanced metadata
        formatted_results = []
        for res in results:
            result_item = {
                "source": res['source'],
                "excerpt": res['text'][:500] + "..." if len(res['text']) > 500 else res['text'],
                "full_content": res['text'],
                "score": round(res.get('score', 0), 4),
                "url": f"http://localhost:8000/static/materials/{res['source']}"
            }
            
            # Add code-specific metadata if available
            if res.get('match_reasons'):
                result_item["match_reasons"] = res['match_reasons']
                result_item["structural_boost"] = res.get('structural_boost', 0)
                result_item["base_score"] = res.get('base_score', 0)
            if res.get('language'):
                result_item["language"] = res['language']
                
            formatted_results.append(result_item)
        
        # Determine search type used
        search_type = "semantic"
        if use_code_aware and is_code_query(request.query):
            search_type = "semantic + syntax-aware"
        
        return {
            "query": request.query,
            "results": formatted_results,
            "total": len(formatted_results),
            "filters": {
                "category": request.category,
                "course": request.course
            },
            "search_type": search_type,
            "is_code_query": is_code_query(request.query)
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

def generate_material_core(prompt: str, category: str, type: str):
    """
    Core logic for generating study materials.
    Returns dict with content, downloadUrl, filename.
    """
    # 1. Retrieve Internal Context (RAG)
    print(f"DEBUG: Retrieving context for generation query: {prompt}")
    context_data = index_manager.retrieve_context(prompt)
    context_str = ""
    if context_data:
        context_str = "## RELEVANT COURSE MATERIALS (Use these as primary source):\n"
        for item in context_data[:5]: # Top 5 sources
            context_str += f"---\nSOURCE: {item['source']}\nCONTENT: {item['text'][:1500]}\n---\n"
    
    # Construct a tailored system prompt
    base_instruction = """
    You are an expert educational content generator for a university-level platform.
    Output MUST be in Markdown format.
    """
    
    specific_instruction = ""
    output_ext = ".md"
    actual_type = type
    
    print(f"DEBUG: generate_material_core called with type='{type}', category='{category}'")
    
    # Handle visual/image generation separately (no LLM content generation needed)
    if type.lower().strip() in ['visual', 'image']:
        print(f"DEBUG: Generating standalone visual for: {prompt}")
        img_path = generate_image(f"Educational diagram or visual aid about: {prompt}", prompt)
        
        if img_path and os.path.exists(img_path):
            img_filename = os.path.basename(img_path)
            return {
                "content": "[IMAGE_GENERATED]",
                "downloadUrl": f"/static/materials/{img_filename}",
                "filename": img_filename,
                "isVisual": True
            }
        else:
            return {
                "content": "Failed to generate visual. Please try again with a more specific description.",
                "downloadUrl": None,
                "filename": None,
                "isVisual": True
            }
    
    if type == 'code':
        if not is_code_request(prompt):
            return {
                "content": "Not a valid question for code generation. Please ask for code specifically.",
                "downloadUrl": None,
                "filename": None
            }
        else:
            lang_name, lang_ext = detect_language_from_prompt(prompt)
            specific_instruction = f"""
            CRITICAL: Return ONLY valid, runnable {lang_name.upper()} code. ABSOLUTELY NO conversational text.
            - Do NOT write the language name at the start.
            - Do NOT say "Here is the code".
            - Do NOT use markdown code blocks (```). Return raw code only.
            - Include comments WITHIN the code to explain logic.
            """
            output_ext = lang_ext
            actual_type = 'code'
    
    elif type == 'reading':
        specific_instruction = """
        Generate CODE-CENTRIC learning material as a structured document:
        - Format with clear H1, H2 headers.
        - Include SHORT code snippets (5-15 lines) to illustrate concepts.
        - Keep explanations concise but thorough.
        """
        output_ext = ".pdf"
        actual_type = 'reading'
    
    elif type == 'slides':
        specific_instruction = """
        PROMPT FOR SLIDES - STRICT FORMATTING RULES:
        1. Use "Slide [N]: [Title]" as the header for each slide.
        2. MAXIMUM 5-6 bullet points per slide. NO MORE.
        3. Each bullet point must be SHORT (max 12 words). Be concise!
        4. Do NOT use nested lists or sub-bullets.
        5. MATH: Use UNICODE symbols (→, ∑, ≤), NOT LaTeX.
        6. Do NOT include conversational text.
        7. Split content across MORE slides if needed - brevity is key!
        8. VISUAL AIDS: Include 2-3 [IMAGE: description] placeholders total.
        """
        output_ext = ".pptx"
        actual_type = 'slides'
    
    elif type == 'pdf' or category == 'theory':
            specific_instruction = """
            Format as a structured document with clear spacing.
            - Use H1 (#) ONLY for the main document title.
            - Use H2 (##) for major sections.
            - Add empty lines between paragraphs.
            """
            output_ext = ".pdf" if type == 'pdf' else ".md"
            actual_type = type

    full_prompt = f"""{base_instruction}
{specific_instruction}

{context_str}

Context: {category} ({actual_type})
User Prompt: {prompt}
"""
    
    CHAT_MODELS = [
        'gemini-2.0-flash',
        'gemini-2.0-flash-lite',
    ]

    
    content = ""
    last_error = None
    
    for model_name in CHAT_MODELS:
        try:
            # Use the Official Google GenAI SDK
            client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
            
            # Tools: Google Search for external context
            tools = [types.Tool(google_search=types.GoogleSearch())]
            
            print(f"DEBUG: Generating content using model={model_name}...")
            response = client.models.generate_content(
                model=model_name,
                contents=full_prompt,
                config=types.GenerateContentConfig(
                    tools=tools,
                    temperature=0.4 # Lower temperature for more grounded content
                )
            )
            content = response.text
            
            # Check for grounding
            if response.candidates and response.candidates[0].grounding_metadata:
                print("DEBUG: Grounding metadata received.")
            
            if content:
                break # Success!

        except Exception as e:
            error_str = str(e)
            print(f"DEBUG: Model {model_name} failed: {error_str[:100]}")
            last_error = e
            # Retry with next model for server issues OR if model doesn't support tools
            if "503" in error_str or "429" in error_str or "overloaded" in error_str.lower() or "unsupported" in error_str.lower():
                continue
            else:
                break # Non-retryable error

    if not content:
        if last_error:
            # Final fallback to a very basic call if all else fails
            try:
                print("DEBUG: Final fallback attempt...")
                client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
                response = client.models.generate_content(
                    model='gemini-2.0-flash-lite', 
                    contents=full_prompt
                )
                content = response.text
            except Exception as e2:
                 raise HTTPException(status_code=500, detail=f"Generation failed completely: {e2}")
        else:
             raise HTTPException(status_code=500, detail="Content generation returned empty response.")

    # Cleanup artifacts
    content = re.sub(r'```\w*\n?', '', content).replace("```", "").strip()

    # File Generation Logic
    filename = f"{actual_type}_{uuid.uuid4().hex[:8]}{output_ext}"
    download_url = None
    
    if actual_type == 'slides':
        create_pptx_from_text(content, filename, topic_context=prompt)
        download_url = f"/static/materials/{filename}"
        content = f"[System] Presentation generated. Download below.\n\nPreview:\n{content[:500]}..."
        
    elif actual_type == 'pdf':
        create_pdf_from_text(content, filename)
        download_url = f"/static/materials/{filename}"
    
    elif actual_type == 'notes':
            filepath = os.path.join(UPLOAD_DIR, filename)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            download_url = f"/static/materials/{filename}"
    
    elif actual_type == 'reading':
            create_pdf_from_text(content, filename)
            download_url = f"/static/materials/{filename}"
    
    if actual_type == 'code':
            lang_name_detected, _ = detect_language_from_prompt(prompt) 
            content = clean_code_content(content, lang_name_detected)
            
            # --- validation start ---
            validation = validate_code_with_fallback(content, lang_name_detected)
            if not validation['valid']:
                print(f"DEBUG: Code validation failed. Error: {validation['error'][:100] if validation.get('error') else 'Unknown error'}...")
                # SELF-CORRECTION: Retry once
                print("DEBUG: Attempting Self-Correction...")
                
                correction_prompt = f"""
The previous code you generated for '{lang_name_detected}' failed to run/compile.
Error Message:
{validation['error']}

Fix the code and return ONLY the valid, runnable code.
"""
                
                retry_prompt = f"{full_prompt}\n\nPREVIOUS_OUTPUT:\n{content}\n\nERROR:\n{correction_prompt}"
                
                try:
                    client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
                    response = client.models.generate_content(
                        model=chosen_model,
                        contents=retry_prompt
                    )
                    fixed_content = response.text
                    content = clean_code_content(fixed_content, lang_name_detected)
                    content = re.sub(r'```\w*\n?', '', content).replace("```", "").strip()
                    print("DEBUG: Self-Correction successful (new content generated).")
                    
                    # Re-validate the fixed code
                    re_validation = validate_code_with_fallback(content, lang_name_detected)
                    if not re_validation['valid']:
                        error_comment = f"\n\n/* \n[System] Code Verification Failed even after correction.\nError: {re_validation['error']}\n*/"
                        if lang_name_detected == 'python':
                            error_comment = f"\n\n'''\n[System] Code Verification Failed even after correction.\nError: {re_validation['error']}\n'''"
                        content += error_comment
                        
                except Exception as e:
                    print(f"DEBUG: Self-Correction failed: {e}")
                    # Append original error
                    error_comment = f"\n\n/* \n[System] Code Verification Failed.\nError: {validation['error']}\n*/"
                    if lang_name_detected == 'python':
                        error_comment = f"\n\n'''\n[System] Code Verification Failed.\nError: {validation['error']}\n'''"
                    content += error_comment

            # --- validation end ---
            
            filepath = os.path.join(UPLOAD_DIR, filename)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            download_url = f"/static/materials/{filename}"

    return {
        "content": content,
        "downloadUrl": download_url,
        "filename": filename
    }

@app.post("/api/generate")
async def generate_content(request: GenerateRequest, user=Depends(optional_current_user)):
    if user:
        repo.log_event(
            event_type="generate",
            uid=user.get("uid"),
            role=user.get("role"),
            metadata={"target": request.prompt[:120], "category": request.category, "type": request.type},
        )
    """
    Generates educational content based on prompt, category, and type.
    Now enhanced with RAG (Internal Context) and Google Search (External Context).
    """
    try:
        return generate_material_core(request.prompt, request.category, request.type)
    except Exception as e:
        print(f"Generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/digitize")
async def digitize_notes(file: UploadFile = File(...), user=Depends(optional_current_user)):
    if user:
        repo.log_event(
            event_type="digitize",
            uid=user.get("uid"),
            role=user.get("role"),
            metadata={"target": file.filename or "handwritten-notes"},
        )
    """
    Bonus: Uses Gemini Vision to transcribe handwritten notes.
    Includes fallback logic for rate limits.
    """
    try:
        content = await file.read()
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
        
        prompt = "Transcribe these handwritten notes into clean, structured Markdown. If there are equations, use LaTeX format."
        
        # List of vision-capable models to try in order
        VISION_MODELS = [
            'gemini-2.0-flash',
            'gemini-2.0-flash-lite',
            'gemini-1.5-flash',
        ]
        
        last_error = None
        
        for model in VISION_MODELS:
            try:
                print(f"DEBUG: Attempting digitization with {model}...")
                response = client.models.generate_content(
                    model=model,
                    contents=[
                        types.Part.from_bytes(data=content, mime_type=file.content_type),
                        prompt
                    ]
                )
                print(f"DEBUG: Success with {model}")
                
                # Generate PDF for the transcription
                filename = f"digitized_{uuid.uuid4().hex[:8]}.pdf"
                create_pdf_from_text(response.text, filename)
                
                return {
                    "transcription": response.text,
                    "downloadUrl": f"/static/materials/{filename}"
                }
                
            except Exception as e:
                print(f"WARN: Model {model} failed: {e}")
                last_error = e
                # Backoff for rate limits
                if "429" in str(e) or "503" in str(e) or "Resource exhausted" in str(e):
                    time.sleep(1)
                    continue
                else:
                    continue

        # If we get here, all models failed
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to digitize notes after trying multiple models. Last error: {str(last_error)}"
        )

    except Exception as e:
        print(f"Error during digitization: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(UPLOAD_DIR, filename)
    if os.path.exists(file_path):
        from fastapi.responses import FileResponse
        return FileResponse(file_path, filename=filename)
    raise HTTPException(status_code=404, detail="File not found")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
